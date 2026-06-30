#!/usr/bin/env python3
"""Remove NotebookLM watermarks from PPTX artifacts.

This is intentionally conservative by default. It removes visible PPTX shapes
that either contain NotebookLM/Google watermark text or look like small
lower-right logo objects when the caller opts into corner-logo removal.

When --inpaint-raster-watermark is enabled, it also attempts an optional
OpenCV-based repair for NotebookLM watermarks baked into ppt/media images.
"""

from __future__ import annotations

import argparse
import json
import re
import shutil
import subprocess
import sys
import time
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from notebooklm_raster_watermark import RasterWatermarkConfig, clean_pptx_media_images


WATERMARK_PATTERNS = [
    r"\bnotebook\s*lm\b",
    r"\bgoogle\s+notebook\s*lm\b",
    r"\bcreated\s+with\s+notebook\s*lm\b",
    r"\bgenerated\s+with\s+notebook\s*lm\b",
    r"\bmade\s+with\s+notebook\s*lm\b",
    r"\bpowered\s+by\s+notebook\s*lm\b",
]

EMU_PER_INCH = 914400


def utc_now() -> str:
    return datetime.now(timezone.utc).replace(microsecond=0).isoformat()


def write_json(path: Path, payload: Any) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(payload, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def read_shape_xml(shape: Any) -> str:
    try:
        return shape._element.xml.lower()
    except Exception:
        return ""


def visible_text(shape: Any) -> str:
    return str(getattr(shape, "text", "") or "").strip()


def shape_label_blob(shape: Any) -> str:
    parts = [
        str(getattr(shape, "name", "") or ""),
        visible_text(shape),
        read_shape_xml(shape),
    ]
    return "\n".join(part for part in parts if part).lower()


def matches_watermark_text(shape: Any, patterns: list[re.Pattern[str]]) -> bool:
    blob = shape_label_blob(shape)
    return any(pattern.search(blob) for pattern in patterns)


def shape_geometry(shape: Any) -> dict[str, int]:
    left = int(getattr(shape, "left", 0) or 0)
    top = int(getattr(shape, "top", 0) or 0)
    width = int(getattr(shape, "width", 0) or 0)
    height = int(getattr(shape, "height", 0) or 0)
    return {
        "left": left,
        "top": top,
        "width": width,
        "height": height,
        "right": left + width,
        "bottom": top + height,
    }


def is_small_lower_right_logo(shape: Any, slide_w: int, slide_h: int) -> bool:
    geo = shape_geometry(shape)
    if geo["width"] <= 0 or geo["height"] <= 0:
        return False
    if getattr(shape, "shape_type", None) not in {
        MSO_SHAPE_TYPE.PICTURE,
        MSO_SHAPE_TYPE.GROUP,
        MSO_SHAPE_TYPE.AUTO_SHAPE,
        MSO_SHAPE_TYPE.TEXT_BOX,
    }:
        return False
    lower_right = geo["right"] >= slide_w * 0.80 and geo["bottom"] >= slide_h * 0.80
    small = geo["width"] <= slide_w * 0.24 and geo["height"] <= slide_h * 0.14
    return lower_right and small


def remove_shape(shape: Any) -> None:
    element = shape._element
    parent = element.getparent()
    if parent is not None:
        parent.remove(element)


def inspect_and_strip_shape_tree(
    shapes: Any,
    *,
    container: str,
    slide_no: int | None,
    slide_w: int,
    slide_h: int,
    patterns: list[re.Pattern[str]],
    remove_corner_logo: bool,
    dry_run: bool,
) -> tuple[list[dict[str, Any]], list[dict[str, Any]]]:
    removed: list[dict[str, Any]] = []
    candidates: list[dict[str, Any]] = []
    for shape in list(shapes):
        text_hit = matches_watermark_text(shape, patterns)
        corner_hit = remove_corner_logo and is_small_lower_right_logo(shape, slide_w, slide_h)
        if not (text_hit or corner_hit):
            continue
        geo = shape_geometry(shape)
        record = {
            "container": container,
            "slide_no": slide_no,
            "shape_name": str(getattr(shape, "name", "") or ""),
            "shape_type": str(getattr(shape, "shape_type", "") or ""),
            "text": visible_text(shape)[:200],
            "reason": "watermark_text" if text_hit else "small_lower_right_logo",
            "geometry_inches": {
                key: round(value / EMU_PER_INCH, 3)
                for key, value in geo.items()
            },
        }
        candidates.append(record)
        if not dry_run:
            remove_shape(shape)
        removed.append(record)
    return removed, candidates


def count_full_slide_raster_slides(presentation: Presentation) -> dict[str, Any]:
    slide_w = int(presentation.slide_width)
    slide_h = int(presentation.slide_height)
    slides: list[dict[str, Any]] = []
    count = 0
    for idx, slide in enumerate(presentation.slides, start=1):
        full = 0
        for shape in slide.shapes:
            if getattr(shape, "shape_type", None) != MSO_SHAPE_TYPE.PICTURE:
                continue
            geo = shape_geometry(shape)
            covers_width = geo["width"] >= slide_w * 0.88 and geo["left"] <= slide_w * 0.06
            covers_height = geo["height"] >= slide_h * 0.88 and geo["top"] <= slide_h * 0.06
            if covers_width and covers_height:
                full += 1
        if full:
            count += 1
        slides.append({"slide_no": idx, "full_slide_picture_count": full})
    return {
        "slide_count": len(slides),
        "full_slide_raster_slide_count": count,
        "full_slide_raster_ratio": count / len(slides) if slides else 0,
        "slides": slides,
    }


def strip_pptx(
    input_pptx: Path,
    output_pptx: Path,
    *,
    remove_corner_logo: bool,
    dry_run: bool,
    inpaint_raster_watermark: bool,
    inpaint_all_media: bool,
    raster_inpaint_min_ratio: float,
    raster_margin_x: int,
    raster_margin_y: int,
    raster_scale: float,
    raster_debug_dir: str,
) -> dict[str, Any]:
    presentation = Presentation(str(input_pptx))
    slide_w = int(presentation.slide_width)
    slide_h = int(presentation.slide_height)
    patterns = [re.compile(pattern, re.I) for pattern in WATERMARK_PATTERNS]
    removed: list[dict[str, Any]] = []
    candidates: list[dict[str, Any]] = []

    for idx, slide in enumerate(presentation.slides, start=1):
        slide_removed, slide_candidates = inspect_and_strip_shape_tree(
            slide.shapes,
            container="slide",
            slide_no=idx,
            slide_w=slide_w,
            slide_h=slide_h,
            patterns=patterns,
            remove_corner_logo=remove_corner_logo,
            dry_run=dry_run,
        )
        removed.extend(slide_removed)
        candidates.extend(slide_candidates)

    for master_idx, master in enumerate(presentation.slide_masters, start=1):
        master_removed, master_candidates = inspect_and_strip_shape_tree(
            master.shapes,
            container=f"slide_master_{master_idx}",
            slide_no=None,
            slide_w=slide_w,
            slide_h=slide_h,
            patterns=patterns,
            remove_corner_logo=remove_corner_logo,
            dry_run=dry_run,
        )
        removed.extend(master_removed)
        candidates.extend(master_candidates)
        for layout_idx, layout in enumerate(master.slide_layouts, start=1):
            layout_removed, layout_candidates = inspect_and_strip_shape_tree(
                layout.shapes,
                container=f"slide_master_{master_idx}/layout_{layout_idx}",
                slide_no=None,
                slide_w=slide_w,
                slide_h=slide_h,
                patterns=patterns,
                remove_corner_logo=remove_corner_logo,
                dry_run=dry_run,
            )
            removed.extend(layout_removed)
            candidates.extend(layout_candidates)

    raster = count_full_slide_raster_slides(presentation)
    if dry_run:
        output_path = ""
    else:
        output_pptx.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(str(output_pptx))
        output_path = str(output_pptx)

    raster_inpaint: dict[str, Any] = {
        "status": "disabled",
        "requested": inpaint_raster_watermark,
        "patched_images": 0,
    }
    if inpaint_raster_watermark:
        if dry_run:
            raster_inpaint = {"status": "skipped_dry_run", "requested": True, "patched_images": 0}
        elif raster["full_slide_raster_ratio"] < raster_inpaint_min_ratio and not inpaint_all_media:
            raster_inpaint = {
                "status": "skipped_low_raster_ratio",
                "requested": True,
                "patched_images": 0,
                "full_slide_raster_ratio": raster["full_slide_raster_ratio"],
                "min_ratio": raster_inpaint_min_ratio,
                "hint": "Use --inpaint-all-media to scan embedded media even when the deck is not mostly full-slide raster images.",
            }
        else:
            raster_inpaint = clean_pptx_media_images(
                output_pptx,
                output_pptx,
                config=RasterWatermarkConfig(
                    search_margin_x=raster_margin_x,
                    search_margin_y=raster_margin_y,
                    scale=raster_scale,
                    debug_dir=raster_debug_dir,
                ),
            )
    warnings: list[str] = []
    if raster["full_slide_raster_ratio"] >= 0.8 and not removed and not inpaint_raster_watermark:
        warnings.append(
            "The deck is mostly full-slide raster images and no editable watermark shapes were found. "
            "Use --inpaint-raster-watermark to attempt optional OpenCV cleanup of baked-in lower-right watermarks."
        )
    if inpaint_raster_watermark:
        status = str(raster_inpaint.get("status", ""))
        if status == "missing_dependency":
            warnings.append(
                "Raster watermark inpainting was requested but optional dependencies are missing: "
                + ", ".join(raster_inpaint.get("missing", []))
            )
        elif status not in {"completed", "skipped_low_raster_ratio", "skipped_dry_run"}:
            warnings.append(f"Raster watermark inpainting did not complete: {status or 'unknown status'}.")
        elif status == "completed" and int(raster_inpaint.get("patched_images", 0) or 0) == 0 and raster["full_slide_raster_ratio"] >= 0.8:
            warnings.append(
                "Raster watermark inpainting scanned the deck but did not patch any embedded images. "
                "The watermark may be absent, too different from the expected NotebookLM mark, or already removed."
            )
    return {
        "schema_version": "1.0.0",
        "tool": "qiaomu-ppt/scripts/strip_notebooklm_watermark.py",
        "generated_at": utc_now(),
        "input": str(input_pptx),
        "output": output_path,
        "dry_run": dry_run,
        "remove_corner_logo": remove_corner_logo,
        "removed_count": len(removed),
        "removed": removed,
        "candidates": candidates,
        "raster_analysis": raster,
        "raster_inpaint": raster_inpaint,
        "warnings": warnings,
        "ok": not warnings or bool(removed) or int(raster_inpaint.get("patched_images", 0) or 0) > 0,
    }


def which_any(names: list[str]) -> str | None:
    for name in names:
        found = shutil.which(name)
        if found:
            return found
    return None


def regenerate_pdf_from_pptx(pptx: Path, output_pdf: Path, *, timeout: int) -> dict[str, Any]:
    soffice = which_any(["soffice", "libreoffice"])
    if not soffice:
        return {"status": "missing", "reason": "LibreOffice/soffice not found"}
    output_pdf.parent.mkdir(parents=True, exist_ok=True)
    before = set(output_pdf.parent.glob("*.pdf"))
    cmd = [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(output_pdf.parent), str(pptx)]
    started = time.time()
    proc = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True, timeout=timeout)
    produced = output_pdf.parent / f"{pptx.stem}.pdf"
    if proc.returncode != 0 or not produced.exists():
        after = set(output_pdf.parent.glob("*.pdf"))
        produced_candidates = sorted(after - before, key=lambda path: path.stat().st_mtime, reverse=True)
        if produced_candidates:
            produced = produced_candidates[0]
        else:
            return {
                "status": "failed",
                "reason": "LibreOffice PDF export failed",
                "command": cmd,
                "stdout": proc.stdout[-2000:],
                "stderr": proc.stderr[-2000:],
                "duration_seconds": round(time.time() - started, 2),
            }
    if produced != output_pdf:
        if output_pdf.exists():
            output_pdf.unlink()
        produced.rename(output_pdf)
    return {
        "status": "exported",
        "path": str(output_pdf),
        "tool": soffice,
        "command": cmd,
        "duration_seconds": round(time.time() - started, 2),
    }


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("input_pptx", type=Path)
    parser.add_argument("--output", type=Path, default=None, help="Cleaned PPTX path. Defaults to <stem>.clean.pptx.")
    parser.add_argument("--report", type=Path, default=None, help="JSON report path.")
    parser.add_argument(
        "--remove-corner-logo",
        action="store_true",
        help="Also remove small lower-right logo/image/text shapes even when they have no useful alt text.",
    )
    parser.add_argument("--dry-run", action="store_true", help="Inspect candidates without writing the output PPTX.")
    parser.add_argument("--pdf-output", type=Path, default=None, help="Regenerate a PDF from the cleaned PPTX.")
    parser.add_argument("--pdf-timeout", type=int, default=180)
    parser.add_argument(
        "--inpaint-raster-watermark",
        action="store_true",
        help="Attempt optional OpenCV cleanup of NotebookLM watermarks baked into ppt/media images.",
    )
    parser.add_argument(
        "--inpaint-all-media",
        action="store_true",
        help="Scan embedded media even when the deck is not mostly full-slide raster images.",
    )
    parser.add_argument(
        "--raster-inpaint-min-ratio",
        type=float,
        default=0.5,
        help="Minimum full-slide raster ratio required before image cleanup runs without --inpaint-all-media.",
    )
    parser.add_argument("--raster-margin-x", type=int, default=400, help="Right-edge search width for raster cleanup.")
    parser.add_argument("--raster-margin-y", type=int, default=120, help="Bottom-edge search height for raster cleanup.")
    parser.add_argument("--raster-scale", type=float, default=3.0, help="Upscale factor for raster watermark detection.")
    parser.add_argument("--raster-debug-dir", default="", help="Optional directory for raster detection debug masks.")
    args = parser.parse_args()

    input_pptx = args.input_pptx.expanduser().resolve()
    if not input_pptx.exists():
        raise SystemExit(f"PPTX not found: {input_pptx}")
    output = args.output.expanduser().resolve() if args.output else input_pptx.with_name(f"{input_pptx.stem}.clean.pptx")
    report = args.report.expanduser().resolve() if args.report else input_pptx.with_suffix(".watermark-report.json")
    payload = strip_pptx(
        input_pptx,
        output,
        remove_corner_logo=args.remove_corner_logo,
        dry_run=args.dry_run,
        inpaint_raster_watermark=args.inpaint_raster_watermark,
        inpaint_all_media=args.inpaint_all_media,
        raster_inpaint_min_ratio=args.raster_inpaint_min_ratio,
        raster_margin_x=args.raster_margin_x,
        raster_margin_y=args.raster_margin_y,
        raster_scale=args.raster_scale,
        raster_debug_dir=args.raster_debug_dir,
    )
    if args.pdf_output and not args.dry_run:
        payload["pdf_export"] = regenerate_pdf_from_pptx(output, args.pdf_output.expanduser().resolve(), timeout=args.pdf_timeout)
    write_json(report, payload)
    print(json.dumps(payload, ensure_ascii=False, indent=2))
    return 0 if payload.get("ok") else 2


if __name__ == "__main__":
    raise SystemExit(main())

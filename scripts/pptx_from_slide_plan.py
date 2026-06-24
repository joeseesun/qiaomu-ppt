#!/usr/bin/env python3
"""Create an editable PPTX from qiaomu-ppt project contracts.

This exporter is intentionally self-contained. It uses python-pptx directly and
does not call any external PPTX skill.
"""

from __future__ import annotations

import argparse
import json
import math
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_W = 13.333333
SLIDE_H = 7.5


def load_json(path: Path) -> Any:
    with path.open("r", encoding="utf-8") as fh:
        return json.load(fh)


def write_json(path: Path, payload: Any) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(payload, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def iter_slides(plan: Any) -> list[dict[str, Any]]:
    if isinstance(plan, list):
        return [item for item in plan if isinstance(item, dict)]
    if isinstance(plan, dict):
        for key in ("slides", "slide_plan", "pages"):
            value = plan.get(key)
            if isinstance(value, list):
                return [item for item in value if isinstance(item, dict)]
    return []


def hex_to_rgb(value: str, fallback: str = "111827") -> RGBColor:
    raw = str(value or fallback).strip().lstrip("#")
    if len(raw) != 6:
        raw = fallback
    try:
        return RGBColor(int(raw[0:2], 16), int(raw[2:4], 16), int(raw[4:6], 16))
    except Exception:
        fb = fallback.lstrip("#")
        return RGBColor(int(fb[0:2], 16), int(fb[2:4], 16), int(fb[4:6], 16))


def rgb_tuple(value: str, fallback: str = "111827") -> tuple[int, int, int]:
    rgb = hex_to_rgb(value, fallback)
    return int(rgb[0]), int(rgb[1]), int(rgb[2])


def mix_hex(a: str, b: str, t: float) -> str:
    ar, ag, ab = rgb_tuple(a)
    br, bg, bb = rgb_tuple(b)
    vals = [round(ar * (1 - t) + br * t), round(ag * (1 - t) + bg * t), round(ab * (1 - t) + bb * t)]
    return "{:02X}{:02X}{:02X}".format(*vals)


def inches(value: float):
    return Inches(value)


def slide_title(slide: dict[str, Any]) -> str:
    return str(slide.get("claim_title") or slide.get("title") or "").strip()


def slide_points(slide: dict[str, Any]) -> list[str]:
    for key in ("content_points", "bullets", "points", "content"):
        value = slide.get(key)
        if isinstance(value, list):
            return [str(item).strip() for item in value if str(item).strip()]
        if isinstance(value, str) and value.strip():
            return [value.strip()]
    return []


def component_plan(slide: dict[str, Any]) -> dict[str, Any]:
    value = slide.get("component_plan")
    if isinstance(value, dict):
        return value
    return {}


def get_palette(spec: dict[str, Any]) -> dict[str, str]:
    palette = spec.get("palette") or spec.get("colors") or {}
    if not isinstance(palette, dict):
        palette = {}
    return {
        "ink": str(palette.get("ink") or palette.get("text") or "#1A1A1A"),
        "muted": str(palette.get("muted") or "#59606A"),
        "paper": str(palette.get("paper") or palette.get("background") or "#F7F3EA"),
        "dark": str(palette.get("dark") or "#111827"),
        "primary": str(palette.get("primary") or "#284B63"),
        "secondary": str(palette.get("secondary") or "#E8DCC5"),
        "accent": str(palette.get("accent") or "#C85A2E"),
        "soft": str(palette.get("soft") or "#EFE8DA"),
        "white": "#FFFFFF",
    }


def text_box(slide, text: str, x: float, y: float, w: float, h: float, size: int, color: str,
             bold: bool = False, align: int | None = None, font: str = "Aptos", line_spacing: float | None = None):
    box = slide.shapes.add_textbox(inches(x), inches(y), inches(w), inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = inches(0)
    frame.margin_right = inches(0)
    frame.margin_top = inches(0.02)
    frame.margin_bottom = inches(0.02)
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    para = frame.paragraphs[0]
    para.text = text
    if align is not None:
        para.alignment = align
    if line_spacing is not None:
        para.line_spacing = line_spacing
    run = para.runs[0] if para.runs else para.add_run()
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = hex_to_rgb(color)
    return box


def bullet_box(slide, points: list[str], x: float, y: float, w: float, h: float, size: int,
               color: str, accent: str, font: str = "Aptos", max_points: int = 4):
    box = slide.shapes.add_textbox(inches(x), inches(y), inches(w), inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = inches(0)
    frame.margin_right = inches(0)
    frame.margin_top = inches(0)
    frame.margin_bottom = inches(0)
    frame.word_wrap = True
    for idx, point in enumerate(points[:max_points]):
        para = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        para.text = point
        para.space_after = Pt(8)
        para.line_spacing = 1.15
        para.level = 0
        run = para.runs[0] if para.runs else para.add_run()
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = hex_to_rgb(color)
        if idx == 0:
            run.font.bold = True
        para.bullet = False
        marker = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            inches(x - 0.18),
            inches(y + idx * 0.58 + 0.12),
            inches(0.08),
            inches(0.08),
        )
        marker.fill.solid()
        marker.fill.fore_color.rgb = hex_to_rgb(accent)
        marker.line.fill.background()
    return box


def add_rect(slide, x: float, y: float, w: float, h: float, fill: str, line: str | None = None,
             radius: bool = False, transparency: int = 0):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, inches(x), inches(y), inches(w), inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(fill)
    if transparency:
        try:
            shape.fill.transparency = transparency
        except Exception:
            pass
    if line:
        shape.line.color.rgb = hex_to_rgb(line)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_line(slide, x1: float, y1: float, x2: float, y2: float, color: str, width: float = 1.0):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, inches(x1), inches(y1), inches(x2), inches(y2))
    line.line.color.rgb = hex_to_rgb(color)
    line.line.width = Pt(width)
    return line


def add_background(slide, palette: dict[str, str], slide_no: int, rhythm: str):
    dark = rhythm == "anchor" or slide_no in {1}
    base = palette["dark"] if dark else palette["paper"]
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(base)
    accent = palette["accent"]
    primary = palette["primary"]
    if dark:
        add_rect(slide, 8.2, -0.2, 5.8, 8.0, mix_hex(primary, "000000", 0.36), None, False)
        for idx in range(5):
            r = 0.55 + idx * 0.22
            orb = slide.shapes.add_shape(MSO_SHAPE.OVAL, inches(10.2 - idx * 0.45), inches(0.62 + idx * 0.62), inches(r), inches(r))
            orb.fill.solid()
            orb.fill.fore_color.rgb = hex_to_rgb(accent if idx % 2 == 0 else primary)
            try:
                orb.fill.transparency = 58 + idx * 4
            except Exception:
                pass
            orb.line.fill.background()
    else:
        add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, base, None, False)
        band = mix_hex(palette["secondary"], "FFFFFF", 0.22)
        add_rect(slide, 0, 0, 0.18, SLIDE_H, accent, None, False)
        add_rect(slide, 9.2, 0, 4.2, SLIDE_H, band, None, False, 8)


def add_footer(slide, slide_no: int, total: int, palette: dict[str, str], dark: bool):
    color = "FFFFFF" if dark else palette["muted"]
    text_box(slide, f"{slide_no:02d} / {total:02d}", 11.75, 7.05, 0.95, 0.22, 8, color, False, PP_ALIGN.RIGHT)


def render_cover(slide, data: dict[str, Any], palette: dict[str, str], total: int):
    title = slide_title(data)
    subtitle = str(data.get("subtitle") or "20页知识型演示：从物理革命到现代世界").strip()
    text_box(slide, title, 0.74, 1.08, 7.2, 1.6, 45, "FFFFFF", True, None, line_spacing=1.1)
    text_box(slide, subtitle, 0.78, 2.86, 5.6, 0.42, 16, "EDE7DB", False)
    text_box(slide, "E = mc²", 0.78, 4.28, 3.0, 0.62, 32, palette["accent"], True)
    text_box(slide, str(data.get("concrete_anchor") or ""), 0.78, 5.02, 6.0, 0.8, 14, "D8D2C8")
    add_footer(slide, 1, total, palette, True)


def render_section(slide, data: dict[str, Any], palette: dict[str, str], slide_no: int, total: int):
    title = slide_title(data)
    text_box(slide, f"{slide_no:02d}", 0.75, 0.78, 1.05, 0.5, 20, palette["accent"], True)
    text_box(slide, title, 0.75, 1.58, 7.8, 1.35, 38, "FFFFFF", True, line_spacing=1.1)
    points = slide_points(data)
    if points:
        text_box(slide, points[0], 0.78, 3.38, 6.5, 0.64, 17, "EDE7DB")
    text_box(slide, str(data.get("concrete_anchor") or ""), 0.78, 5.7, 7.1, 0.54, 13, "D8D2C8")
    add_footer(slide, slide_no, total, palette, True)


def render_timeline(slide, data: dict[str, Any], palette: dict[str, str], slide_no: int, total: int):
    title = slide_title(data)
    text_box(slide, title, 0.68, 0.52, 8.6, 0.82, 27, palette["ink"], True, line_spacing=1.1)
    points = slide_points(data)
    milestones = component_plan(data).get("milestones")
    if not isinstance(milestones, list) or not milestones:
        milestones = points[:5]
    x0, x1, y = 1.1, 12.0, 3.4
    add_line(slide, x0, y, x1, y, palette["primary"], 1.5)
    for idx, item in enumerate(milestones[:5]):
        x = x0 + (x1 - x0) * idx / max(1, min(len(milestones), 5) - 1)
        marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, inches(x - 0.16), inches(y - 0.16), inches(0.32), inches(0.32))
        marker.fill.solid()
        marker.fill.fore_color.rgb = hex_to_rgb(palette["accent"])
        marker.line.color.rgb = hex_to_rgb(palette["paper"])
        label = str(item.get("label") if isinstance(item, dict) else item)
        year = str(item.get("year") if isinstance(item, dict) else "").strip()
        text_box(slide, year or label.split(" ")[0], x - 0.45, y - 0.75, 0.9, 0.28, 12, palette["primary"], True, PP_ALIGN.CENTER)
        text_box(slide, label, x - 0.9, y + 0.42, 1.8, 0.72, 10, palette["ink"], False, PP_ALIGN.CENTER)
    add_footer(slide, slide_no, total, palette, False)


def render_comparison(slide, data: dict[str, Any], palette: dict[str, str], slide_no: int, total: int):
    title = slide_title(data)
    text_box(slide, title, 0.68, 0.52, 8.8, 0.85, 27, palette["ink"], True, line_spacing=1.1)
    points = slide_points(data)
    left = points[:2] or ["旧图景：时间、空间和物质被看作独立舞台"]
    right = points[2:4] or ["新图景：测量、能量和几何本身进入物理定律"]
    add_rect(slide, 0.78, 1.78, 5.65, 4.72, "FFFFFF", mix_hex(palette["secondary"], "FFFFFF", 0.1), True)
    add_rect(slide, 6.9, 1.78, 5.65, 4.72, "FFFFFF", mix_hex(palette["secondary"], "FFFFFF", 0.1), True)
    text_box(slide, "之前的问题", 1.12, 2.08, 2.4, 0.36, 16, palette["primary"], True)
    text_box(slide, "爱因斯坦的转向", 7.24, 2.08, 2.7, 0.36, 16, palette["accent"], True)
    bullet_box(slide, left, 1.22, 2.86, 4.72, 1.75, 15, palette["ink"], palette["primary"], max_points=2)
    bullet_box(slide, right, 7.34, 2.86, 4.72, 1.75, 15, palette["ink"], palette["accent"], max_points=2)
    add_footer(slide, slide_no, total, palette, False)


def render_equation(slide, data: dict[str, Any], palette: dict[str, str], slide_no: int, total: int):
    title = slide_title(data)
    text_box(slide, title, 0.68, 0.52, 8.9, 0.85, 27, palette["ink"], True, line_spacing=1.1)
    equation = str(component_plan(data).get("equation") or data.get("equation") or "E = mc²")
    add_rect(slide, 0.88, 1.92, 5.15, 2.1, palette["dark"], None, True)
    text_box(slide, equation, 1.2, 2.48, 4.45, 0.74, 42, "FFFFFF", True, PP_ALIGN.CENTER)
    points = slide_points(data)
    bullet_box(slide, points[:4], 6.72, 1.96, 5.25, 3.0, 15, palette["ink"], palette["accent"], max_points=4)
    text_box(slide, str(data.get("concrete_anchor") or ""), 0.92, 4.78, 10.8, 0.52, 12, palette["muted"])
    add_footer(slide, slide_no, total, palette, False)


def render_cards(slide, data: dict[str, Any], palette: dict[str, str], slide_no: int, total: int):
    title = slide_title(data)
    text_box(slide, title, 0.68, 0.52, 8.8, 0.85, 27, palette["ink"], True, line_spacing=1.1)
    points = slide_points(data)[:4]
    positions = [(0.82, 1.72), (6.83, 1.72), (0.82, 4.08), (6.83, 4.08)]
    for idx, (point, (x, y)) in enumerate(zip(points, positions), start=1):
        add_rect(slide, x, y, 5.42, 1.72, "FFFFFF", mix_hex(palette["secondary"], "FFFFFF", 0.18), True)
        text_box(slide, f"{idx:02d}", x + 0.28, y + 0.23, 0.55, 0.26, 11, palette["accent"], True)
        text_box(slide, point, x + 0.28, y + 0.6, 4.72, 0.72, 14, palette["ink"], True, line_spacing=1.12)
    add_footer(slide, slide_no, total, palette, False)


def render_model(slide, data: dict[str, Any], palette: dict[str, str], slide_no: int, total: int):
    title = slide_title(data)
    text_box(slide, title, 0.68, 0.52, 8.7, 0.82, 27, palette["ink"], True, line_spacing=1.1)
    labels = component_plan(data).get("nodes")
    if not isinstance(labels, list) or not labels:
        labels = slide_points(data)[:4]
    labels = [str(item) for item in labels[:4]]
    while len(labels) < 4:
        labels.append("思想实验")
    center_x, center_y = 6.6, 3.75
    add_rect(slide, center_x - 1.1, center_y - 0.43, 2.2, 0.86, palette["dark"], None, True)
    text_box(slide, str(component_plan(data).get("center") or "物理图景"), center_x - 0.92, center_y - 0.16, 1.85, 0.24, 12, "FFFFFF", True, PP_ALIGN.CENTER)
    positions = [(2.0, 2.05), (10.0, 2.05), (2.0, 5.15), (10.0, 5.15)]
    for label, (x, y) in zip(labels, positions):
        add_line(slide, center_x, center_y, x + 0.75, y + 0.25, palette["primary"], 0.9)
        add_rect(slide, x, y, 1.5, 0.5, palette["soft"], None, True)
        text_box(slide, label, x + 0.13, y + 0.13, 1.24, 0.2, 10, palette["ink"], True, PP_ALIGN.CENTER)
    add_footer(slide, slide_no, total, palette, False)


def render_quote(slide, data: dict[str, Any], palette: dict[str, str], slide_no: int, total: int):
    title = slide_title(data)
    quote = str(component_plan(data).get("quote") or (slide_points(data)[0] if slide_points(data) else title))
    text_box(slide, "“", 0.76, 0.72, 0.8, 0.8, 54, palette["accent"], True)
    text_box(slide, title, 1.28, 0.9, 8.3, 0.7, 25, "FFFFFF", True, line_spacing=1.1)
    text_box(slide, quote, 1.34, 2.42, 8.4, 1.6, 28, "FFFFFF", True, line_spacing=1.14)
    text_box(slide, str(data.get("concrete_anchor") or ""), 1.36, 5.16, 7.5, 0.46, 13, "E0D8CC")
    add_footer(slide, slide_no, total, palette, True)


def render_default(slide, data: dict[str, Any], palette: dict[str, str], slide_no: int, total: int):
    title = slide_title(data)
    text_box(slide, title, 0.68, 0.52, 8.8, 0.85, 27, palette["ink"], True, line_spacing=1.1)
    points = slide_points(data)[:4]
    bullet_box(slide, points, 0.96, 1.76, 5.8, 3.0, 15, palette["ink"], palette["accent"], max_points=4)
    add_rect(slide, 8.0, 1.62, 3.2, 3.2, palette["dark"], None, True)
    text_box(slide, str(data.get("visual_role") or "Proof"), 8.28, 2.78, 2.65, 0.32, 16, "FFFFFF", True, PP_ALIGN.CENTER)
    text_box(slide, str(data.get("concrete_anchor") or ""), 0.96, 5.72, 8.8, 0.5, 12, palette["muted"])
    add_footer(slide, slide_no, total, palette, False)


def render_slide(prs: Presentation, data: dict[str, Any], palette: dict[str, str], slide_no: int, total: int):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    rhythm = str(data.get("rhythm") or ("anchor" if slide_no in {1, total} else "dense")).lower()
    add_background(slide, palette, slide_no, rhythm)
    component = str(component_plan(data).get("component_type") or data.get("component_type") or "").lower()
    layout = str(data.get("layout_pattern") or data.get("layout_pattern_id") or "").lower()
    dark = rhythm == "anchor" or slide_no == 1
    if slide_no == 1:
        render_cover(slide, data, palette, total)
    elif rhythm == "anchor" and ("quote" in component or "quote" in layout):
        render_quote(slide, data, palette, slide_no, total)
    elif rhythm == "anchor":
        render_section(slide, data, palette, slide_no, total)
    elif "timeline" in component or "timeline" in layout:
        render_timeline(slide, data, palette, slide_no, total)
    elif "comparison" in component or "compare" in component:
        render_comparison(slide, data, palette, slide_no, total)
    elif "formula" in component or "equation" in component:
        render_equation(slide, data, palette, slide_no, total)
    elif "model" in component or "loop" in component or "map" in component:
        render_model(slide, data, palette, slide_no, total)
    elif "cards" in component or "grid" in component:
        render_cards(slide, data, palette, slide_no, total)
    else:
        render_default(slide, data, palette, slide_no, total)
    if not dark:
        # Tiny semantic motif, intentionally foreground and editable.
        x = 12.82
        y = 0.32 + (slide_no % 5) * 0.16
        mark = slide.shapes.add_shape(MSO_SHAPE.OVAL, inches(x), inches(y), inches(0.16), inches(0.16))
        mark.fill.solid()
        mark.fill.fore_color.rgb = hex_to_rgb(palette["accent"])
        mark.line.fill.background()
    return slide


def build_pptx(project: Path, output: Path) -> dict[str, Any]:
    plan_path = project / "slide_plan.json"
    spec_path = project / "spec_lock.json"
    if not plan_path.exists():
        raise SystemExit(f"missing slide_plan.json: {plan_path}")
    slides = iter_slides(load_json(plan_path))
    if not slides:
        raise SystemExit("slide_plan.json has no slides")
    spec = load_json(spec_path) if spec_path.exists() else {}
    if not isinstance(spec, dict):
        spec = {}
    palette = get_palette(spec)

    prs = Presentation()
    prs.slide_width = inches(SLIDE_W)
    prs.slide_height = inches(SLIDE_H)
    for idx, slide_data in enumerate(slides, start=1):
        render_slide(prs, slide_data, palette, idx, len(slides))
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    manifest = {
        "schema_version": "1.0.0",
        "generator": "qiaomu-ppt/scripts/pptx_from_slide_plan.py",
        "runtime_dependency": "python-pptx",
        "external_skill_dependency": "none",
        "project": str(project),
        "output": str(output),
        "slide_count": len(slides),
        "source_contracts": {
            "slide_plan": "slide_plan.json",
            "spec_lock": "spec_lock.json" if spec_path.exists() else "",
            "visual_contract": "visual_contract.json" if (project / "visual_contract.json").exists() else "",
        },
        "editable_policy": "Text, shapes, connectors, and diagrams are authored as editable PowerPoint objects where python-pptx supports them.",
    }
    write_json(project / "pptx_generation_manifest.json", manifest)
    return manifest


def main() -> int:
    parser = argparse.ArgumentParser(description="Create an editable PPTX from qiaomu-ppt slide_plan/spec_lock contracts.")
    parser.add_argument("project", help="qiaomu-ppt project directory")
    parser.add_argument("--output", "-o", help="Output PPTX path. Defaults to <project>/exports/<project>.pptx")
    args = parser.parse_args()
    project = Path(args.project).resolve()
    output = Path(args.output).resolve() if args.output else project / "exports" / f"{project.name}.pptx"
    manifest = build_pptx(project, output)
    print(json.dumps(manifest, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

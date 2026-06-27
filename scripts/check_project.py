#!/usr/bin/env python3
"""Check qiaomu-ppt project artifacts without depending on upstream skills."""

from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path
from typing import Any


REQUIRED_ROOT = ["deck_brief.md", "slide_plan.json", "style_brief.md"]
SLIDE_FIELDS = [
    "slide_no",
    "title",
    "intent",
    "concrete_anchor",
    "visual_role",
    "speaker_note_goal",
    "qa_risk",
]
PLACEHOLDERS = [
    r"\[必填\]",
    r"SLIDES_HERE",
    r"lorem ipsum",
    r"\bTODO\b",
    r"\bTBD\b",
]
SVG_BANNED = [
    r"<style\b",
    r"\bclass=",
    r"<foreignObject\b",
    r"<mask\b",
    r"<symbol\b",
    r"<use\b",
    r"textPath",
    r"@font-face",
    r"<animate",
    r"<script\b",
    r"\bon[a-zA-Z]+=",
    r"<iframe\b",
]
VISUAL_BACKGROUND_ROLES_MIN = 4
VISUAL_BACKGROUND_MAX_REPEAT = 2
ALLOWED_VISUAL_NOISE_BUDGETS = {"quiet", "moderate", "expressive"}
DEFAULT_MAX_ACTIVE_COLORS = 3
DEFAULT_MAX_CARDLIKE_CONTAINERS = 4
DEFAULT_MAX_COVER_METRIC_CHIPS = 3
DEFAULT_MIN_BACKGROUND_ASSETS_LONG_DECK = 8
DEFAULT_UNIQUE_BACKGROUND_TARGET = 12
DEFAULT_MAX_PICTURE_ASPECT_DISTORTION = 0.03
MIN_CJK_TITLE_LINE_HEIGHT = 1.03
WARN_CJK_TITLE_LINE_HEIGHT = 1.08
REQUIRED_TYPOGRAPHY_POLICY_FIELDS = {
    "primary_title_dominant",
    "secondary_text_max_ratio",
    "promote_oversized_support_to_title",
}
REQUIRED_BACKGROUND_FORBIDDEN_OBJECTS = {
    "box",
    "card",
    "panel",
    "frame",
    "placeholder",
    "chart area",
    "image slot",
    "ui chrome",
    "text block",
}
REQUIRED_EVIDENCE_LAYOUT_TREATMENTS = {
    "full_chart",
    "chart_with_takeaway",
    "chart_crop",
    "chart_then_takeaway",
}
REQUIRED_SHAPE_COMPONENT_POLICY_FIELDS = {
    "safe_area_policy",
    "text_fit_policy",
    "connector_policy",
    "operator_policy",
    "card_density_policy",
    "separator_policy",
    "preview_rejection_policy",
}
MAINLINE_SKIP_INTENT_TOKENS = {
    "cover",
    "chapter",
    "section",
    "divider",
    "closing",
    "summary",
    "navigation",
    "breathing",
    "封面",
    "章节",
    "分隔",
    "收束",
    "结尾",
    "目录",
}
IMAGE_SLOT_FIELDS = [
    "slot_id",
    "slide_no",
    "x",
    "y",
    "w",
    "h",
    "fit",
    "mask",
    "padding",
    "overflow_policy",
]
LAYOUT_CONTRACT_SLIDE_FIELDS = [
    "rhythm",
    "proof_object",
    "layout_pattern_id",
    "component_type",
    "reading_path",
    "coordinate_slots",
    "group_ids",
]
COORDINATE_SLOT_FIELDS = ["slot_id", "x", "y", "w", "h"]
ALLOWED_PAGE_RHYTHMS = {"anchor", "dense", "breathing"}
PROOF_SLOT_TOKENS = {
    "proof",
    "chart",
    "table",
    "diagram",
    "process",
    "timeline",
    "media",
    "image",
    "hero",
    "quote",
    "map",
    "screenshot",
    "comparison",
    "model",
    "number",
    "kpi",
}
SVG_DIR_NAMES = {"svg", "svgs", "svg_output", "svg-final", "svg_final", "slides"}
SVG_EXCLUDED_ROOT_DIRS = {"previews", "_preview_work", "assets", "sources", "html", "html-parity", "exports"}
CONTENT_CONTRACT_FIELDS = [
    "audience",
    "purpose",
    "desired_action",
    "current_state",
    "desired_state",
    "stakes",
    "structure_framework",
    "title_policy",
    "copy_density",
    "evidence_policy",
    "speaker_note_policy",
    "slide_claims",
]
ALLOWED_STRUCTURE_FRAMEWORKS = {
    "pyramid",
    "scqa",
    "mece",
    "storyline",
    "teaching_arc",
    "hybrid",
}
WEAK_TITLE_LABELS = {
    "agenda",
    "overview",
    "background",
    "problem",
    "solution",
    "data",
    "summary",
    "conclusion",
    "目录",
    "概览",
    "背景",
    "现状",
    "问题",
    "方案",
    "数据",
    "总结",
    "结论",
}
VISIBLE_INTERNAL_METADATA = [
    ("fetched via", r"\bfetched\s+via\b"),
    ("generated with", r"\bgenerated\s+with\b"),
    ("qiaomu-markdown-proxy", r"qiaomu-markdown-proxy"),
    ("speaker cue", r"Speaker\s+cue\s*:"),
]
VISIBLE_PRODUCTION_JARGON = [
    ("deck", r"\bdeck\b"),
    ("route", r"\broute\b"),
    ("fallback", r"\bfallback\b"),
    ("artifact", r"\bartifact\b"),
    ("pipeline", r"\bpipeline\b"),
    ("source_fetch", r"\bsource_fetch\b"),
    ("pptx export", r"\bPPTX\s+export\b"),
    ("路线卡", r"路线卡"),
    ("制作方法", r"制作方法|工作流|方法论|质量门|检查报告"),
    ("版式术语", r"版式|布局模式|\bL\d{2}\b|\bITL\d{2}\b"),
    ("生成术语", r"生成策略|生图策略|图片生成|生成图|提示词|负面提示"),
    ("交付术语", r"交付格式|可编辑\s*PPTX|导出状态"),
    ("图版标签", r"(?:复原|化石|骨骼|证据)?图版[：:]?"),
    ("讲解顺序", r"讲解顺序[：:]?"),
    ("证据编号栏", r"证据\s*\d{1,2}\s*/"),
]
VISUAL_ASSET_REQUIRED_FIELDS = [
    "asset_id",
    "filename",
    "path",
    "purpose",
    "asset_role",
    "acquire_via",
    "status",
    "reference",
    "page_role",
    "text_policy",
    "aspect_ratio",
    "editable_policy",
]
VISUAL_ASSET_ACQUIRE_VIA = {"ai", "web", "user", "source", "formula", "placeholder"}
VISUAL_ASSET_TERMINAL_STATUS = {
    "ai": {"Generated", "Needs-Manual", "Missing", "Failed"},
    "web": {"Sourced", "Needs-Manual", "Missing", "Failed"},
    "user": {"Existing", "Needs-Manual", "Missing"},
    "source": {"Existing", "Needs-Manual", "Missing"},
    "formula": {"Rendered", "Needs-Manual", "Missing"},
    "placeholder": {"Placeholder", "Needs-Manual", "Missing"},
}
VISUAL_ASSET_FILE_STATUSES = {"Generated", "Sourced", "Existing", "Rendered"}
URL_PATTERN = re.compile(r"https?://[^\s)>\"]+", flags=re.IGNORECASE)
DATA_URI_PATTERN = re.compile(r"data:[^\"')>\s]+", flags=re.IGNORECASE)
FORMAL_HTML_SCREENSHOT_PATTERNS = [
    r"rendered_slide_parity",
    r"html_parity_preview",
    r"html_from_previews",
    r"<img[^>]+(?:previews/|render/page|slide-\d{2}\.(?:jpg|jpeg|png))",
]
TITLE_SELECTOR_TOKENS = (
    "h1",
    "h2",
    "h3",
    "title",
    "headline",
    "claim",
    "cover-sub",
    "closing",
)


def read_text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


def load_json(path: Path) -> Any:
    with path.open("r", encoding="utf-8") as fh:
        return json.load(fh)


def artifact_mtime(path: Path) -> float:
    if path.is_dir():
        newest = path.stat().st_mtime
        for child in path.rglob("*"):
            try:
                newest = max(newest, child.stat().st_mtime)
            except OSError:
                continue
        return newest
    return path.stat().st_mtime


def iter_slides(plan: Any) -> list[dict[str, Any]]:
    if isinstance(plan, list):
        return [item for item in plan if isinstance(item, dict)]
    if isinstance(plan, dict):
        for key in ("slides", "slide_plan", "pages"):
            value = plan.get(key)
            if isinstance(value, list):
                return [item for item in value if isinstance(item, dict)]
    return []


def iter_image_slots(contract: Any) -> list[dict[str, Any]]:
    if not isinstance(contract, dict):
        return []
    slots = contract.get("image_slots")
    if isinstance(slots, list):
        return [slot for slot in slots if isinstance(slot, dict)]
    return []


def iter_slide_palette_slots(contract: Any) -> list[dict[str, Any]]:
    if not isinstance(contract, dict):
        return []
    slots = contract.get("slide_palette_slots")
    if isinstance(slots, list):
        return [slot for slot in slots if isinstance(slot, dict)]
    return []


def iter_slide_claims(contract: Any) -> list[dict[str, Any]]:
    if not isinstance(contract, dict):
        return []
    claims = contract.get("slide_claims")
    if isinstance(claims, list):
        return [claim for claim in claims if isinstance(claim, dict)]
    return []


def normalize_frameworks(value: Any) -> list[str]:
    if isinstance(value, str):
        return [value.lower()]
    if isinstance(value, list):
        return [str(item).lower() for item in value]
    return []


def slide_title(slide: dict[str, Any]) -> str:
    return str(slide.get("claim_title") or slide.get("title") or "").strip()


def slide_points(slide: dict[str, Any]) -> list[Any]:
    for key in ("bullets", "content_points", "points", "content"):
        value = slide.get(key)
        if isinstance(value, list):
            return value
    return []


def slide_number_value(slide: dict[str, Any]) -> int | None:
    value = slide.get("slide_no") or slide.get("page") or slide.get("slide")
    try:
        return int(value)
    except Exception:
        return None


def slide_identifier_values(slide: dict[str, Any]) -> set[str]:
    values: set[str] = set()
    number = slide_number_value(slide)
    if number is not None:
        values.add(str(number))
        values.add(f"{number:02d}")
        values.add(f"p{number:02d}")
        values.add(f"slide{number:02d}")
        values.add(f"slide-{number:02d}")
        values.add(f"slide_{number:02d}")
    for key in ("slide_id", "id", "page_id"):
        value = str(slide.get(key) or "").strip().lower()
        if value:
            values.add(value)
    return values


def layout_pattern_ids() -> set[str]:
    library = Path(__file__).resolve().parents[1] / "references" / "layout-pattern-library.md"
    if not library.exists():
        return set()
    text = read_text(library)
    return {match.group(0).upper() for match in re.finditer(r"\bL\d{2}\b", text)}


def has_svg_pages(root: Path) -> bool:
    return any(is_formal_svg_page(root, path) for path in root.rglob("*.svg"))


def is_formal_svg_page(root: Path, path: Path) -> bool:
    if not path.is_file() or path.suffix.lower() != ".svg":
        return False
    parts = path.relative_to(root).parts
    parent_parts = parts[:-1]
    if not parent_parts:
        return False
    if parent_parts[0] in SVG_EXCLUDED_ROOT_DIRS:
        return False
    return any(part in SVG_DIR_NAMES for part in parent_parts)


def group_ids_in_svg(path: Path) -> set[str]:
    text = read_text(path)
    ids: set[str] = set()
    for match in re.finditer(r"<g\b[^>]*\bid\s*=\s*(['\"])(.*?)\1", text, flags=re.IGNORECASE | re.DOTALL):
        group_id = match.group(2).strip()
        if group_id:
            ids.add(group_id)
    return ids


def element_ids_in_svg(path: Path) -> set[str]:
    text = read_text(path)
    ids: set[str] = set()
    for match in re.finditer(r"\bid\s*=\s*(['\"])(.*?)\1", text, flags=re.IGNORECASE | re.DOTALL):
        element_id = match.group(2).strip()
        if element_id:
            ids.add(element_id)
    return ids


def top_level_group_count(path: Path) -> int:
    text = read_text(path)
    open_groups = 0
    count = 0
    for match in re.finditer(r"</?g\b[^>]*>|<g\b[^>]*/>", text, flags=re.IGNORECASE | re.DOTALL):
        token = match.group(0)
        is_close = token.startswith("</")
        is_self_close = token.endswith("/>")
        if is_close:
            open_groups = max(0, open_groups - 1)
            continue
        if open_groups == 0:
            count += 1
        if not is_self_close:
            open_groups += 1
    return count


def svg_candidates_for_slide(root: Path, slide: dict[str, Any]) -> list[Path]:
    number = slide_number_value(slide)
    custom_identifiers = {
        str(slide.get(key) or "").strip().lower()
        for key in ("slide_id", "id", "page_id")
        if str(slide.get(key) or "").strip()
    }
    candidates: list[Path] = []
    for path in sorted(root.rglob("*.svg")):
        if not is_formal_svg_page(root, path):
            continue
        stem = path.stem.lower()
        compact = re.sub(r"[^a-z0-9]+", "", stem)
        matched = False
        if number is not None:
            numeric_patterns = [
                rf"(?<!\d){re.escape(str(number))}(?!\d)",
                rf"(?<!\d){number:02d}(?!\d)",
            ]
            matched = any(
                re.search(pattern, target)
                for pattern in numeric_patterns
                for target in (stem, compact)
            )
        if not matched:
            matched = any(
                identifier and (identifier in stem or identifier in compact)
                for identifier in custom_identifiers
            )
        if matched:
            candidates.append(path)
    return candidates


def all_svg_group_ids(root: Path) -> dict[str, set[str]]:
    result: dict[str, set[str]] = {}
    for path in sorted(root.rglob("*.svg")):
        if not is_formal_svg_page(root, path):
            continue
        result[path.relative_to(root).as_posix()] = group_ids_in_svg(path)
    return result


def normalize_contract_slides(contract: Any) -> list[dict[str, Any]]:
    if not isinstance(contract, dict):
        return []
    slides = contract.get("slides") or contract.get("pages")
    if isinstance(slides, list):
        return [slide for slide in slides if isinstance(slide, dict)]
    if isinstance(slides, dict):
        normalized: list[dict[str, Any]] = []
        for key, value in slides.items():
            if isinstance(value, dict):
                record = dict(value)
                record.setdefault("slide_id", key)
                normalized.append(record)
        return normalized
    return []


def animation_targets(animation_value: Any) -> list[tuple[str, list[str]]]:
    targets: list[tuple[str, list[str]]] = []
    if not isinstance(animation_value, dict):
        return targets
    slides = animation_value.get("slides")
    if isinstance(slides, dict):
        for slide_key, slide_animation in slides.items():
            groups: list[str] = []
            if isinstance(slide_animation, dict):
                raw_groups = slide_animation.get("groups")
                if isinstance(raw_groups, dict):
                    groups.extend(str(group_id) for group_id in raw_groups.keys())
                elif isinstance(raw_groups, list):
                    for item in raw_groups:
                        if isinstance(item, str):
                            groups.append(item)
                        elif isinstance(item, dict):
                            for key in ("group_id", "target", "id"):
                                value = str(item.get(key) or "").strip()
                                if value:
                                    groups.append(value)
                for key in ("group_id", "target", "id"):
                    value = str(slide_animation.get(key) or "").strip()
                    if value:
                        groups.append(value)
            targets.append((str(slide_key), sorted(set(groups))))
    elif isinstance(slides, list):
        for idx, slide_animation in enumerate(slides, start=1):
            if not isinstance(slide_animation, dict):
                continue
            slide_key = str(
                slide_animation.get("slide_id")
                or slide_animation.get("slide_no")
                or slide_animation.get("page")
                or idx
            )
            groups = []
            raw_groups = slide_animation.get("groups")
            if isinstance(raw_groups, dict):
                groups.extend(str(group_id) for group_id in raw_groups.keys())
            elif isinstance(raw_groups, list):
                for item in raw_groups:
                    if isinstance(item, str):
                        groups.append(item)
                    elif isinstance(item, dict):
                        value = str(item.get("group_id") or item.get("target") or item.get("id") or "").strip()
                        if value:
                            groups.append(value)
            targets.append((slide_key, sorted(set(groups))))
    return targets


def is_weak_title(title: str) -> bool:
    normalized = re.sub(r"[\s:：\-—–_]+", "", title.strip().lower())
    if normalized in WEAK_TITLE_LABELS:
        return True
    if title.strip().lower() in WEAK_TITLE_LABELS:
        return True
    return False


def check_content_contract(path: Path, slides: list[dict[str, Any]]) -> tuple[list[str], list[str], dict[str, Any]]:
    failures: list[str] = []
    warnings: list[str] = []
    evidence: dict[str, Any] = {}
    try:
        contract = load_json(path)
    except Exception as exc:
        return [f"invalid content_contract.json: {exc}"], warnings, evidence

    if not isinstance(contract, dict):
        return ["content_contract.json must be an object"], warnings, evidence

    missing = [field for field in CONTENT_CONTRACT_FIELDS if not contract.get(field)]
    if missing:
        failures.append(f"content_contract.json missing fields: {', '.join(missing)}")

    frameworks = normalize_frameworks(contract.get("structure_framework"))
    evidence["structure_framework"] = frameworks
    unknown = [item for item in frameworks if item not in ALLOWED_STRUCTURE_FRAMEWORKS]
    if unknown:
        warnings.append(f"content_contract.json has unknown structure framework: {', '.join(unknown)}")

    title_policy = str(contract.get("title_policy") or "").lower()
    if "claim" not in title_policy:
        warnings.append("content_contract.json should use claim_titles as title_policy")

    claims = iter_slide_claims(contract)
    evidence["slide_claim_count"] = len(claims)
    if slides and len(claims) < len(slides):
        failures.append("content_contract.json needs slide_claims for every slide")

    weak_titles: list[str] = []
    dense_slides: list[str] = []
    for slide in slides:
        title = slide_title(slide)
        slide_no = slide.get("slide_no", "?")
        if not title:
            failures.append(f"slide {slide_no} has no title or claim_title")
        elif is_weak_title(title):
            weak_titles.append(f"slide {slide_no}: {title}")
        if re.search(r"(\.\.\.|…)", title):
            failures.append(f"slide {slide_no} title appears truncated with ellipsis: {title}")
        compact_title = re.sub(r"\s+", "", title)
        if re.search(r"[\u4e00-\u9fff]", title) and len(compact_title) > 54:
            failures.append(f"slide {slide_no} CJK title is too long for reliable PPTX/HTML rendering: {title}")
        elif len(title) > 96:
            failures.append(f"slide {slide_no} title is too long for reliable PPTX/HTML rendering: {title}")
        points = slide_points(slide)
        if len(points) > 5:
            dense_slides.append(f"slide {slide_no}: {len(points)} visible chunks")

    if weak_titles:
        failures.append(
            "slides use weak label titles instead of claim titles: " + "; ".join(weak_titles[:8])
        )
    if dense_slides:
        warnings.append(
            "slides may exceed 3-5 visible support chunks: " + "; ".join(dense_slides[:8])
        )

    return failures, warnings, evidence


def truthy(value: Any) -> bool:
    if isinstance(value, bool):
        return value
    if isinstance(value, str):
        return value.strip().lower() in {"1", "true", "yes", "y", "required"}
    return bool(value)


def is_skipped_research(contract: dict[str, Any]) -> bool:
    status = str(contract.get("research_status") or "").strip().lower()
    return status in {"skipped", "skipped_by_user", "user_skipped"}


def is_mainline_slide(slide: dict[str, Any], index: int, total: int) -> bool:
    joined = " ".join(
        str(slide.get(key) or "").lower()
        for key in ("intent", "visual_role", "rhythm", "layout_pattern", "component_type", "title", "claim_title")
    )
    if any(token in joined for token in MAINLINE_SKIP_INTENT_TOKENS):
        return False
    if total > 1 and index in {1, total}:
        role = str(slide.get("visual_role") or slide.get("intent") or "").lower()
        if any(token in role for token in ("cover", "closing", "封面", "结尾")):
            return False
    return True


def normalize_id_list(value: Any) -> list[str]:
    if isinstance(value, str):
        return [value] if value.strip() else []
    if isinstance(value, list):
        return [str(item).strip() for item in value if str(item).strip()]
    return []


def check_topic_research_artifacts(
    root: Path,
    slides: list[dict[str, Any]],
    content_contract: dict[str, Any] | None,
) -> tuple[list[str], list[str], dict[str, Any]]:
    failures: list[str] = []
    warnings: list[str] = []
    evidence: dict[str, Any] = {}
    if not isinstance(content_contract, dict) or not truthy(content_contract.get("research_required")):
        return failures, warnings, evidence

    evidence["research_required"] = True
    research_status = str(content_contract.get("research_status") or "").strip()
    evidence["research_status"] = research_status
    skipped = is_skipped_research(content_contract)

    for field in ("research_status", "topic_angle", "source_coverage", "evidence_policy"):
        if not content_contract.get(field):
            failures.append(f"content_contract.json research_required=true missing {field}")

    research_plan_paths = [root / "research_plan.md", root / "research_plan.json"]
    if not any(path.exists() for path in research_plan_paths):
        message = "topic-researched deck needs research_plan.md or research_plan.json"
        (warnings if skipped else failures).append(message)
    evidence["research_plan"] = next((path.name for path in research_plan_paths if path.exists()), "")

    research_dossier_path = root / "research_dossier.md"
    evidence["research_dossier"] = (
        str(research_dossier_path.relative_to(root))
        if research_dossier_path.exists()
        else ""
    )
    if not research_dossier_path.exists() and not skipped:
        warnings.append(
            "topic-researched deck should include research_dossier.md as the reviewable Markdown synthesis before slide planning"
        )

    sources_dir = root / "sources"
    required_source_files = ["source_manifest.json", "source_notes.md", "source_cards.json"]
    for name in required_source_files:
        path = sources_dir / name
        if not path.exists():
            message = f"topic-researched deck needs sources/{name}"
            (warnings if skipped else failures).append(message)

    source_cards_path = sources_dir / "source_cards.json"
    source_card_ids: set[str] = set()
    if source_cards_path.exists():
        try:
            source_cards_payload = load_json(source_cards_path)
        except Exception as exc:
            failures.append(f"invalid sources/source_cards.json: {exc}")
            source_cards_payload = None
        if isinstance(source_cards_payload, dict):
            cards = source_cards_payload.get("cards")
            if not isinstance(cards, list) or not cards:
                failures.append("sources/source_cards.json needs a non-empty cards list")
            else:
                for idx, card in enumerate(cards, start=1):
                    if not isinstance(card, dict):
                        failures.append(f"source card {idx} must be an object")
                        continue
                    card_id = str(card.get("id") or "").strip()
                    if not card_id:
                        failures.append(f"source card {idx} missing id")
                        continue
                    source_card_ids.add(card_id)
                    for field in ("claim", "evidence", "confidence"):
                        if not card.get(field):
                            warnings.append(f"source card {card_id} missing {field}")
            image_candidates = source_cards_payload.get("image_candidates")
            evidence["image_candidate_count"] = len(image_candidates) if isinstance(image_candidates, list) else 0
            if not image_candidates:
                warnings.append("topic-researched deck has no image_candidates; record source image gaps or generated-image policy")
        else:
            failures.append("sources/source_cards.json must be an object")
    evidence["source_card_count"] = len(source_card_ids)

    if slides:
        missing_slide_cards: list[str] = []
        unknown_slide_cards: list[str] = []
        total = len(slides)
        for idx, slide in enumerate(slides, start=1):
            if not is_mainline_slide(slide, idx, total):
                continue
            slide_no = slide.get("slide_no") or slide.get("page") or idx
            ids = normalize_id_list(slide.get("source_card_ids"))
            if not ids:
                missing_slide_cards.append(str(slide_no))
                continue
            if source_card_ids:
                unknown = [card_id for card_id in ids if card_id not in source_card_ids]
                if unknown:
                    unknown_slide_cards.append(f"{slide_no}: {', '.join(unknown)}")
            if not slide.get("source_anchor"):
                warnings.append(f"slide {slide_no} has source_card_ids but no source_anchor")
        if missing_slide_cards:
            failures.append(
                "topic-researched mainline slides missing source_card_ids: "
                + ", ".join(missing_slide_cards[:12])
            )
        if unknown_slide_cards:
            failures.append(
                "topic-researched slides reference unknown source_card_ids: "
                + "; ".join(unknown_slide_cards[:8])
            )

    return failures, warnings, evidence


def check_visual_contract(path: Path, slide_count: int) -> tuple[list[str], list[str], dict[str, Any]]:
    failures: list[str] = []
    warnings: list[str] = []
    evidence: dict[str, Any] = {}
    try:
        contract = load_json(path)
    except Exception as exc:
        return [f"invalid visual_contract.json: {exc}"], warnings, evidence
    if not isinstance(contract, dict):
        return ["visual_contract.json must be an object"], warnings, evidence

    visual_noise_budget = contract.get("visual_noise_budget")
    evidence["visual_noise_budget"] = visual_noise_budget
    if slide_count > 8 and not visual_noise_budget:
        failures.append("visual_contract.json needs visual_noise_budget for decks over 8 slides")
    if visual_noise_budget and str(visual_noise_budget) not in ALLOWED_VISUAL_NOISE_BUDGETS:
        failures.append(
            "visual_contract.json visual_noise_budget must be one of: "
            + ", ".join(sorted(ALLOWED_VISUAL_NOISE_BUDGETS))
        )
    if visual_noise_budget and str(visual_noise_budget) != "quiet":
        warnings.append("visual_noise_budget is not quiet; verify this is intentional and content-led")

    color_budget = contract.get("color_budget")
    max_active_colors = DEFAULT_MAX_ACTIVE_COLORS
    if slide_count > 8 and not isinstance(color_budget, dict):
        failures.append("visual_contract.json needs color_budget for decks over 8 slides")
    if isinstance(color_budget, dict):
        declared_max = color_budget.get("max_active_colors_per_slide")
        try:
            max_active_colors = int(declared_max)
        except Exception:
            failures.append("color_budget.max_active_colors_per_slide must be an integer")
        if max_active_colors > DEFAULT_MAX_ACTIVE_COLORS:
            failures.append(
                f"color_budget.max_active_colors_per_slide must be <= {DEFAULT_MAX_ACTIVE_COLORS}"
            )
        accent_policy = str(color_budget.get("accent_policy") or color_budget.get("default_formula") or "").lower()
        if "one" not in accent_policy and "1" not in accent_policy:
            warnings.append("color_budget should declare a one-accent-per-slide policy")
    evidence["max_active_colors_per_slide"] = max_active_colors

    palette_slots = iter_slide_palette_slots(contract)
    evidence["slide_palette_slot_count"] = len(palette_slots)
    if slide_count > 8 and len(palette_slots) < slide_count:
        failures.append("visual_contract.json needs slide_palette_slots for every slide")
    for idx, slot in enumerate(palette_slots, start=1):
        colors = slot.get("active_colors")
        if not isinstance(colors, list) or not colors:
            failures.append(f"slide_palette_slot {idx} needs non-empty active_colors")
            continue
        unique_colors = {str(color).strip().lower() for color in colors if str(color).strip()}
        if len(unique_colors) > max_active_colors:
            failures.append(
                f"slide {slot.get('slide_no', idx)} uses {len(unique_colors)} active colors; max is {max_active_colors}"
            )

    background_policy = contract.get("background_asset_policy")
    if slide_count > 8 and not isinstance(background_policy, dict):
        failures.append("visual_contract.json needs background_asset_policy for decks over 8 slides")
    if isinstance(background_policy, dict):
        mode = str(background_policy.get("mode") or background_policy.get("generated_backgrounds") or "")
        evidence["background_asset_policy"] = mode
        line_policy = str(background_policy.get("decorative_line_policy") or "").lower()
        if "functional" not in line_policy and "forbid" not in line_policy and "no decorative" not in line_policy:
            failures.append(
                "background_asset_policy.decorative_line_policy must forbid non-functional decorative lines"
            )
        atmosphere_policy = str(background_policy.get("atmosphere_only_policy") or "").lower()
        editable_policy = str(background_policy.get("editable_foreground_policy") or "").lower()
        forbidden_objects = background_policy.get("forbidden_generated_objects")
        evidence["background_atmosphere_only"] = bool(atmosphere_policy)
        if "atmosphere" not in atmosphere_policy or "only" not in atmosphere_policy:
            failures.append(
                "background_asset_policy.atmosphere_only_policy must declare that generated backgrounds are atmosphere-only"
            )
        procedural_policy = str(background_policy.get("procedural_fallback_policy") or "").lower()
        evidence["procedural_fallback_policy"] = procedural_policy
        if slide_count > 8 and not procedural_policy:
            failures.append(
                "background_asset_policy.procedural_fallback_policy must declare CSS/Canvas/SVG fallback when image generation is unavailable"
            )
        if procedural_policy and not any(term in procedural_policy for term in ("css", "canvas", "svg", "procedural")):
            warnings.append(
                "background_asset_policy.procedural_fallback_policy should mention procedural CSS, Canvas, or SVG fallback"
            )
        if "editable" not in editable_policy or "foreground" not in editable_policy:
            failures.append(
                "background_asset_policy.editable_foreground_policy must declare that layout objects stay editable foreground objects"
            )
        if not isinstance(forbidden_objects, list) or not forbidden_objects:
            failures.append("background_asset_policy.forbidden_generated_objects must list forbidden baked-in layout objects")
        else:
            normalized_forbidden = {
                str(item).strip().lower()
                for item in forbidden_objects
                if str(item).strip()
            }
            missing_forbidden = sorted(REQUIRED_BACKGROUND_FORBIDDEN_OBJECTS - normalized_forbidden)
            if missing_forbidden:
                failures.append(
                    "background_asset_policy.forbidden_generated_objects missing: "
                    + ", ".join(missing_forbidden)
                )

    roles = contract.get("background_roles")
    if not isinstance(roles, list) or len(roles) < VISUAL_BACKGROUND_ROLES_MIN:
        failures.append(
            f"visual_contract.json needs at least {VISUAL_BACKGROUND_ROLES_MIN} background_roles"
        )
    evidence["background_role_count"] = len(roles) if isinstance(roles, list) else 0

    slide_roles = contract.get("slide_roles")
    if slide_count > 0:
        if not isinstance(slide_roles, list) or len(slide_roles) < slide_count:
            failures.append("visual_contract.json needs per-slide slide_roles for every slide")
        else:
            previous = None
            run_length = 0
            used_roles: set[str] = set()
            used_assets: set[str] = set()
            for idx, item in enumerate(slide_roles, start=1):
                missing_role_fields = [
                    field
                    for field in ("slide_no", "layout_role", "reading_path", "background_role", "background_asset", "dominant_object")
                    if not item.get(field)
                ]
                if missing_role_fields:
                    failures.append(
                        f"slide_role {idx} missing fields: {', '.join(missing_role_fields)}"
                    )
                role = str(item.get("background_role") or "")
                if role:
                    used_roles.add(role)
                asset = str(item.get("background_asset") or "").strip()
                if asset:
                    used_assets.add(asset)
                if role and role == previous:
                    run_length += 1
                else:
                    previous = role
                    run_length = 1
                if role and run_length > VISUAL_BACKGROUND_MAX_REPEAT:
                    failures.append(
                        f"background_role `{role}` repeats more than {VISUAL_BACKGROUND_MAX_REPEAT} consecutive slides"
                    )
                    break
            if slide_count > 8 and len(used_roles) < VISUAL_BACKGROUND_ROLES_MIN:
                failures.append(
                    f"deck has {slide_count} slides but uses only {len(used_roles)} background roles"
                )
            unique_background_target = min(slide_count, DEFAULT_UNIQUE_BACKGROUND_TARGET)
            evidence["unique_background_asset_count"] = len(used_assets)
            if slide_count > 8 and len(used_assets) < unique_background_target:
                failures.append(
                    f"deck has {slide_count} slides but uses only {len(used_assets)} unique background assets; target is {unique_background_target}"
                )
            previous_asset = None
            for idx, item in enumerate(slide_roles, start=1):
                asset = str(item.get("background_asset") or "").strip()
                if asset and previous_asset and asset == previous_asset:
                    failures.append(
                        f"slide_role {idx} reuses the exact same background_asset as the previous slide: {asset}"
                    )
                    break
                if asset:
                    previous_asset = asset

    if isinstance(background_policy, dict):
        paths = background_policy.get("background_paths") or background_policy.get("procedural_background_paths")
        if isinstance(paths, dict):
            evidence["background_asset_count"] = len(paths)
            if slide_count > 8 and len(paths) < DEFAULT_MIN_BACKGROUND_ASSETS_LONG_DECK:
                warnings.append(
                    f"deck has {slide_count} slides but only {len(paths)} declared background assets; prefer {DEFAULT_MIN_BACKGROUND_ASSETS_LONG_DECK}-12 variants"
                )

    layout_policy = contract.get("layout_quality_policy")
    if slide_count > 8 and not isinstance(layout_policy, dict):
        failures.append("visual_contract.json needs layout_quality_policy for decks over 8 slides")
    if isinstance(layout_policy, dict):
        evidence["layout_quality_policy"] = layout_policy.get("composition_formula", "")
        filler_policy = str(layout_policy.get("decorative_filler_policy") or "").lower()
        if "forbid" not in filler_policy and "no" not in filler_policy:
            failures.append("layout_quality_policy.decorative_filler_policy must forbid decorative filler")
        try:
            max_regions = int(layout_policy.get("max_primary_regions", 0))
        except Exception:
            max_regions = 0
        if max_regions <= 0 or max_regions > 3:
            failures.append("layout_quality_policy.max_primary_regions must be between 1 and 3")

    typography_policy = contract.get("typography_hierarchy_policy")
    if slide_count > 8 and not isinstance(typography_policy, dict):
        failures.append("visual_contract.json needs typography_hierarchy_policy for decks over 8 slides")
    if isinstance(typography_policy, dict):
        evidence["typography_hierarchy_policy"] = typography_policy
        missing_typography = [
            field
            for field in sorted(REQUIRED_TYPOGRAPHY_POLICY_FIELDS)
            if field not in typography_policy
        ]
        if missing_typography:
            failures.append(
                "typography_hierarchy_policy missing fields: "
                + ", ".join(missing_typography)
            )
        if typography_policy.get("primary_title_dominant") is not True:
            failures.append("typography_hierarchy_policy.primary_title_dominant must be true")
        try:
            ratio = float(typography_policy.get("secondary_text_max_ratio", 0))
        except Exception:
            ratio = 0
        if ratio <= 0 or ratio > 1:
            failures.append("typography_hierarchy_policy.secondary_text_max_ratio must be > 0 and <= 1")
        promote_policy = str(typography_policy.get("promote_oversized_support_to_title") or "").lower()
        if "promote" not in promote_policy and "title" not in promote_policy:
            failures.append(
                "typography_hierarchy_policy.promote_oversized_support_to_title must require promoting oversized support text to title"
            )
        title_line_policy = str(typography_policy.get("title_line_height_policy") or "").lower()
        evidence["title_line_height_policy"] = title_line_policy
        if slide_count > 8 and not title_line_policy:
            warnings.append("typography_hierarchy_policy should declare title_line_height_policy for CJK multi-line titles")
        elif title_line_policy and not any(term in title_line_policy for term in ("line", "height", "leading", "行高")):
            warnings.append("typography_hierarchy_policy.title_line_height_policy should explicitly mention line-height/leading/行高")

    chrome_policy = contract.get("layout_chrome_policy")
    if slide_count > 8 and not isinstance(chrome_policy, dict):
        failures.append("visual_contract.json needs layout_chrome_policy for decks over 8 slides")
    if isinstance(chrome_policy, dict):
        evidence["layout_chrome_policy"] = chrome_policy
        nested_policy = str(chrome_policy.get("nested_card_policy") or "").lower()
        parent_policy = str(chrome_policy.get("parent_container_policy") or "").lower()
        bottom_row_policy = str(chrome_policy.get("bottom_card_row_policy") or "").lower()
        badge_policy = str(chrome_policy.get("numbered_badge_policy") or "").lower()
        if "no" not in nested_policy and "forbid" not in nested_policy:
            failures.append("layout_chrome_policy.nested_card_policy must forbid nested visible cards")
        if "no" not in parent_policy and "forbid" not in parent_policy:
            failures.append("layout_chrome_policy.parent_container_policy must forbid giant parent containers used only for grouping")
        if "no" not in bottom_row_policy and "process" not in bottom_row_policy and "timeline" not in bottom_row_policy:
            failures.append("layout_chrome_policy.bottom_card_row_policy must restrict bottom card rows to real process/timeline use")
        if "sequence" not in badge_policy and "ordered" not in badge_policy:
            failures.append("layout_chrome_policy.numbered_badge_policy must restrict numbered badges to real sequences")
        try:
            max_cardlike = int(chrome_policy.get("max_cardlike_containers_per_slide", 0))
        except Exception:
            max_cardlike = 0
        if max_cardlike <= 0 or max_cardlike > DEFAULT_MAX_CARDLIKE_CONTAINERS:
            failures.append(
                f"layout_chrome_policy.max_cardlike_containers_per_slide must be 1-{DEFAULT_MAX_CARDLIKE_CONTAINERS}"
            )
        try:
            cover_chips = int(chrome_policy.get("max_cover_metric_chips", 0))
        except Exception:
            cover_chips = 0
        if cover_chips <= 0 or cover_chips > DEFAULT_MAX_COVER_METRIC_CHIPS:
            failures.append(
                f"layout_chrome_policy.max_cover_metric_chips must be 1-{DEFAULT_MAX_COVER_METRIC_CHIPS}"
            )

    shape_policy = contract.get("shape_component_policy")
    if slide_count > 8 and not isinstance(shape_policy, dict):
        failures.append("visual_contract.json needs shape_component_policy for decks over 8 slides")
    if isinstance(shape_policy, dict):
        evidence["shape_component_policy"] = shape_policy
        missing_shape_policy = [
            field
            for field in sorted(REQUIRED_SHAPE_COMPONENT_POLICY_FIELDS)
            if not shape_policy.get(field)
        ]
        if missing_shape_policy:
            failures.append(
                "shape_component_policy missing fields: "
                + ", ".join(missing_shape_policy)
            )
        safe_area_policy = str(shape_policy.get("safe_area_policy") or "").lower()
        if not any(term in safe_area_policy for term in ("margin", "padding", "safe")):
            failures.append("shape_component_policy.safe_area_policy must require safe margins/padding")
        text_fit_policy = str(shape_policy.get("text_fit_policy") or "").lower()
        if not any(term in text_fit_policy for term in ("inside", "within", "fit")):
            failures.append("shape_component_policy.text_fit_policy must require text to fit inside visible shapes")
        if not any(term in text_fit_policy for term in ("fail", "split", "enlarge", "shorten", "resize")):
            failures.append("shape_component_policy.text_fit_policy must say how to fix or fail overflowing text")
        connector_policy = str(shape_policy.get("connector_policy") or "").lower()
        if not any(term in connector_policy for term in ("thin", "simple", "line", "arrowhead", "whitespace")):
            failures.append("shape_component_policy.connector_policy must prefer thin/simple connectors")
        if not any(term in connector_policy for term in ("chevron", "block", "chunky")):
            warnings.append("shape_component_policy.connector_policy should explicitly restrict chevrons/block arrows")
        if not any(term in connector_policy for term in ("perimeter", "edge", "port", "boundary")):
            failures.append("shape_component_policy.connector_policy must require perimeter/edge/port connector endpoints")
        if not any(term in connector_policy for term in ("through", "cross", "interior", "inside node", "text")):
            failures.append("shape_component_policy.connector_policy must forbid connectors crossing node interiors or text")
        if not any(term in connector_policy for term in ("align", "grid", "centerline", "axis", "spacing")):
            failures.append("shape_component_policy.connector_policy must require diagram alignment geometry such as grid, axis, centerline, or equal spacing")
        operator_policy = str(shape_policy.get("operator_policy") or "").lower()
        if not any(term in operator_policy for term in ("outside", "standalone", "not inside", "never inside")):
            failures.append("shape_component_policy.operator_policy must keep operators outside arrow shapes")
        card_density_policy = str(shape_policy.get("card_density_policy") or "").lower()
        if not any(term in card_density_policy for term in ("fewer", "limit", "only", "avoid", "max")):
            failures.append("shape_component_policy.card_density_policy must limit card/pill overuse")
        separator_policy = str(shape_policy.get("separator_policy") or "").lower()
        if not any(term in separator_policy for term in ("functional", "meaningful", "zone", "separator")):
            failures.append("shape_component_policy.separator_policy must restrict separators to meaningful zone boundaries")
        preview_policy = str(shape_policy.get("preview_rejection_policy") or "").lower()
        if not any(term in preview_policy for term in ("reject", "fail", "fix")):
            failures.append("shape_component_policy.preview_rejection_policy must reject/fix rendered overflow and connector clutter")

    evidence_policy = contract.get("evidence_layout_policy")
    if slide_count > 8 and not isinstance(evidence_policy, dict):
        failures.append("visual_contract.json needs evidence_layout_policy for decks over 8 slides")
    if isinstance(evidence_policy, dict):
        evidence["evidence_layout_policy"] = evidence_policy
        treatments = evidence_policy.get("allowed_treatments")
        if not isinstance(treatments, list):
            failures.append("evidence_layout_policy.allowed_treatments must list allowed evidence layouts")
        else:
            normalized_treatments = {
                str(item).strip().lower()
                for item in treatments
                if str(item).strip()
            }
            missing_treatments = sorted(REQUIRED_EVIDENCE_LAYOUT_TREATMENTS - normalized_treatments)
            if missing_treatments:
                failures.append(
                    "evidence_layout_policy.allowed_treatments missing: "
                    + ", ".join(missing_treatments)
                )
        try:
            max_same = int(evidence_policy.get("max_consecutive_same_composition", 0))
        except Exception:
            max_same = 0
        if max_same <= 0 or max_same > 2:
            failures.append("evidence_layout_policy.max_consecutive_same_composition must be 1 or 2")
        right_side_policy = str(evidence_policy.get("right_side_policy") or "").lower()
        if "one" not in right_side_policy and "1" not in right_side_policy:
            failures.append("evidence_layout_policy.right_side_policy must require one clear takeaway/number")
        rail_policy = str(evidence_policy.get("rail_policy") or "").lower()
        if "timeline" not in rail_policy and "process" not in rail_policy and "forbid" not in rail_policy:
            failures.append("evidence_layout_policy.rail_policy must forbid generic rails outside timeline/process slides")

    dense_chart_policy = str(contract.get("dense_chart_policy") or "").lower()
    evidence["dense_chart_policy"] = dense_chart_policy
    if slide_count > 8 and not dense_chart_policy:
        failures.append("visual_contract.json needs dense_chart_policy for decks over 8 slides")
    if dense_chart_policy and not any(term in dense_chart_policy for term in ("crop", "split", "zoom", "notes")):
        warnings.append("dense_chart_policy should mention crop, split, zoom, or notes for unreadable dense charts")

    slots = iter_image_slots(contract)
    evidence["image_slot_count"] = len(slots)
    for idx, slot in enumerate(slots, start=1):
        missing = [field for field in IMAGE_SLOT_FIELDS if field not in slot or slot.get(field) in ("", None)]
        if missing:
            failures.append(f"image slot {idx} missing fields: {', '.join(missing)}")
        if slot.get("overflow_policy") != "clip_or_fail":
            warnings.append(f"image slot {idx} should use overflow_policy: clip_or_fail")
        if slot.get("mask") == "rounded_rect" and not slot.get("clip_method"):
            warnings.append(
                f"image slot {idx} uses rounded_rect mask; include clip_method to prove real clipping/compositing"
            )
    return failures, warnings, evidence


def check_visual_asset_manifest(
    root: Path,
    require_terminal: bool,
    *,
    require_real_imagegen: bool = False,
) -> tuple[list[str], list[str], dict[str, Any]]:
    failures: list[str] = []
    warnings: list[str] = []
    evidence: dict[str, Any] = {}
    manifest_path = root / "visual_asset_manifest.json"
    evidence["visual_asset_manifest"] = str(manifest_path.relative_to(root)) if manifest_path.exists() else ""
    if not manifest_path.exists():
        return failures, warnings, evidence

    try:
        manifest = load_json(manifest_path)
    except Exception as exc:
        return [f"invalid visual_asset_manifest.json: {exc}"], warnings, evidence
    if not isinstance(manifest, dict):
        return ["visual_asset_manifest.json must be an object"], warnings, evidence

    model = manifest.get("deck_image_model")
    if not isinstance(model, dict):
        failures.append("visual_asset_manifest.json missing deck_image_model")
        model = {}
    else:
        for field in ("image_rendering", "image_palette_behavior", "color_scheme"):
            if not model.get(field):
                failures.append(f"visual_asset_manifest.deck_image_model missing {field}")
    evidence["deck_image_model"] = {
        "image_rendering": model.get("image_rendering"),
        "image_palette_behavior": model.get("image_palette_behavior"),
    }

    items = manifest.get("items")
    if not isinstance(items, list) or not items:
        failures.append("visual_asset_manifest.json needs a non-empty items list")
        return failures, warnings, evidence

    counts: dict[str, int] = {}
    statuses: dict[str, int] = {}
    ai_count = 0
    ai_generated_count = 0
    procedural_fallback_count = 0
    real_imagegen_count = 0
    needs_manual: list[str] = []
    for idx, item in enumerate(items, start=1):
        if not isinstance(item, dict):
            failures.append(f"visual asset item {idx} must be an object")
            continue
        missing = [field for field in VISUAL_ASSET_REQUIRED_FIELDS if not item.get(field)]
        if missing:
            failures.append(f"visual asset item {idx} missing fields: {', '.join(missing)}")

        acquire_via = str(item.get("acquire_via") or "").strip().lower()
        status = str(item.get("status") or "").strip()
        counts[acquire_via] = counts.get(acquire_via, 0) + 1
        statuses[status] = statuses.get(status, 0) + 1
        if acquire_via not in VISUAL_ASSET_ACQUIRE_VIA:
            failures.append(f"visual asset item {idx} has unsupported acquire_via: {acquire_via}")
            continue
        allowed_statuses = VISUAL_ASSET_TERMINAL_STATUS[acquire_via] | {"Pending"}
        if status not in allowed_statuses:
            failures.append(
                f"visual asset item {idx} has invalid status `{status}` for acquire_via `{acquire_via}`"
            )
        if require_terminal and status == "Pending":
            failures.append(
                f"visual asset item {idx} remains Pending even though export artifacts exist"
            )
        if status == "Needs-Manual":
            needs_manual.append(str(item.get("filename") or item.get("asset_id") or idx))

        if acquire_via == "ai":
            ai_count += 1
            prompt = str(item.get("prompt") or "")
            if not prompt:
                failures.append(f"AI visual asset item {idx} missing prompt")
            elif len(prompt.split()) < 35:
                warnings.append(
                    f"AI visual asset item {idx} prompt may be too short to encode rendering, palette, composition, and hard rules"
                )
            if item.get("page_role") not in {"hero_page", "local"}:
                failures.append(f"AI visual asset item {idx} page_role must be hero_page or local")
            if item.get("text_policy") not in {"none", "embedded"}:
                failures.append(f"AI visual asset item {idx} text_policy must be none or embedded")
            if item.get("page_role") == "local" and not item.get("visual_type"):
                failures.append(f"AI local visual asset item {idx} missing visual_type")
            if item.get("page_role") == "hero_page" and not item.get("hero_primitive"):
                warnings.append(f"AI hero_page visual asset item {idx} should declare hero_primitive")
            art_brief = item.get("art_direction_brief")
            if not isinstance(art_brief, dict):
                warnings.append(f"AI visual asset item {idx} should include art_direction_brief")
            else:
                for field in ("art_direction", "composition", "safe_area", "foreground_boundary", "negative_prompt"):
                    if not art_brief.get(field):
                        warnings.append(f"AI visual asset item {idx} art_direction_brief missing {field}")
            if not item.get("safe_area"):
                warnings.append(f"AI visual asset item {idx} should declare safe_area for editable foreground placement")
            if not item.get("negative_prompt"):
                warnings.append(f"AI visual asset item {idx} should declare negative_prompt")
            if status == "Generated":
                ai_generated_count += 1
                if not item.get("generator"):
                    failures.append(f"AI visual asset item {idx} is Generated but missing generator")
                if not item.get("dimensions"):
                    warnings.append(f"AI visual asset item {idx} is Generated but missing dimensions")
                generator = str(item.get("generator") or "")
                if generator == "procedural-preview-fallback":
                    procedural_fallback_count += 1
                elif generator:
                    real_imagegen_count += 1

        if acquire_via == "web" and status == "Sourced" and not item.get("license_tier"):
            failures.append(f"web visual asset item {idx} is Sourced but missing license_tier")

        if status in VISUAL_ASSET_FILE_STATUSES:
            asset_path = root / str(item.get("path") or "")
            if not asset_path.exists():
                failures.append(
                    f"visual asset item {idx} status `{status}` requires file: {item.get('path')}"
                )

    evidence["visual_asset_counts"] = counts
    evidence["visual_asset_statuses"] = statuses
    evidence["ai_generated_count"] = ai_generated_count
    evidence["real_imagegen_count"] = real_imagegen_count
    evidence["procedural_fallback_count"] = procedural_fallback_count
    evidence["needs_manual_assets"] = needs_manual
    if needs_manual:
        warnings.append("visual assets still need manual action: " + ", ".join(needs_manual[:8]))

    if ai_count:
        direction_path = root / "image_art_direction.json"
        queue_path = root / "assets" / "images" / "image_generation_queue.json"
        evidence["image_art_direction"] = str(direction_path.relative_to(root)) if direction_path.exists() else ""
        evidence["image_generation_queue"] = str(queue_path.relative_to(root)) if queue_path.exists() else ""
        if not direction_path.exists():
            warnings.append("AI visual assets should include image_art_direction.json")
        else:
            try:
                direction = load_json(direction_path)
                direction_items = direction.get("items") if isinstance(direction, dict) else None
                if not isinstance(direction_items, list) or len(direction_items) < ai_count:
                    warnings.append("image_art_direction.json should cover every AI visual asset")
            except Exception as exc:
                failures.append(f"invalid image_art_direction.json: {exc}")
        if not queue_path.exists():
            warnings.append("AI visual assets should include assets/images/image_generation_queue.json")
        else:
            try:
                queue = load_json(queue_path)
                queue_items = queue.get("items") if isinstance(queue, dict) else None
                if not isinstance(queue_items, list) or len(queue_items) < ai_count:
                    warnings.append("assets/images/image_generation_queue.json should cover every AI visual asset")
            except Exception as exc:
                failures.append(f"invalid assets/images/image_generation_queue.json: {exc}")
        prompts_path = root / "assets" / "images" / "image_prompts.json"
        prompts_md = root / "assets" / "images" / "image_prompts.md"
        if not prompts_path.exists():
            failures.append("AI visual assets require assets/images/image_prompts.json")
        else:
            try:
                prompts = load_json(prompts_path)
                prompt_items = prompts.get("items") if isinstance(prompts, dict) else None
                if not isinstance(prompt_items, list) or len(prompt_items) < ai_count:
                    failures.append("assets/images/image_prompts.json does not cover all AI visual assets")
            except Exception as exc:
                failures.append(f"invalid assets/images/image_prompts.json: {exc}")
        if not prompts_md.exists():
            warnings.append("AI visual assets should also render assets/images/image_prompts.md for manual fallback")
        if procedural_fallback_count:
            message = (
                f"{procedural_fallback_count} AI visual asset(s) use procedural-preview-fallback; "
                "replace with real image-generation outputs before claiming final visual quality"
            )
            if require_real_imagegen:
                failures.append(message)
            else:
                warnings.append(message)
        if require_real_imagegen and ai_generated_count and not real_imagegen_count:
            failures.append("strict image generation check requires at least one non-procedural generated AI asset")

    return failures, warnings, evidence


def check_source_image_usage(root: Path, asset_evidence: dict[str, Any]) -> tuple[list[str], dict[str, Any]]:
    warnings: list[str] = []
    evidence: dict[str, Any] = {}
    cards_path = root / "sources" / "source_cards.json"
    if not cards_path.exists():
        return warnings, evidence
    try:
        payload = load_json(cards_path)
    except Exception:
        return warnings, evidence
    image_candidates = payload.get("image_candidates") if isinstance(payload, dict) else []
    candidate_count = len(image_candidates) if isinstance(image_candidates, list) else 0
    evidence["source_image_candidate_count"] = candidate_count
    counts = asset_evidence.get("visual_asset_counts") if isinstance(asset_evidence, dict) else {}
    if not isinstance(counts, dict):
        counts = {}
    source_asset_count = sum(int(counts.get(key) or 0) for key in ("source", "web", "user"))
    evidence["source_backed_visual_asset_count"] = source_asset_count
    if candidate_count and source_asset_count == 0:
        warnings.append(
            "sources/source_cards.json has image_candidates, but visual_asset_manifest.json has no source/web/user visual assets; "
            "use source images where relevant before falling back to AI/procedural visuals"
        )
    return warnings, evidence


def check_source_manifest(root: Path, deck_brief_text: str) -> tuple[list[str], list[str], dict[str, Any]]:
    failures: list[str] = []
    warnings: list[str] = []
    evidence: dict[str, Any] = {}
    urls = URL_PATTERN.findall(deck_brief_text)
    sources_dir = root / "sources"
    manifest_path = sources_dir / "source_manifest.json"
    evidence["deck_brief_url_count"] = len(urls)
    evidence["source_manifest"] = str(manifest_path.relative_to(root)) if manifest_path.exists() else ""

    if urls and not manifest_path.exists():
        failures.append("deck_brief.md references URLs but sources/source_manifest.json is missing")
        return failures, warnings, evidence
    if not manifest_path.exists():
        if sources_dir.exists() and list(sources_dir.glob("*.md")):
            warnings.append("sources/ contains Markdown files but no source_manifest.json")
        return failures, warnings, evidence

    try:
        manifest = load_json(manifest_path)
    except Exception as exc:
        return [f"invalid sources/source_manifest.json: {exc}"], warnings, evidence

    for companion in ("source_notes.md", "source_cards.json"):
        if not (sources_dir / companion).exists():
            warnings.append(f"sources/source_manifest.json exists but sources/{companion} is missing")

    source_records = manifest.get("sources")
    if not isinstance(source_records, list) or not source_records:
        failures.append("sources/source_manifest.json needs a non-empty sources list")
        return failures, warnings, evidence

    evidence["source_count"] = len(source_records)
    for idx, item in enumerate(source_records, start=1):
        if not isinstance(item, dict):
            failures.append(f"source_manifest source {idx} must be an object")
            continue
        for field in ("input", "title", "source_type", "fetch_route", "markdown_path"):
            if not item.get(field):
                failures.append(f"source_manifest source {idx} missing {field}")
        markdown_path = item.get("markdown_path")
        if markdown_path and not (sources_dir / str(markdown_path)).exists():
            failures.append(f"source_manifest source {idx} markdown_path does not exist: {markdown_path}")
        if item.get("missing_evidence"):
            warnings.append(f"source_manifest source {idx} has missing evidence: {item.get('missing_evidence')}")
    return failures, warnings, evidence


def check_layout_execution_contract(root: Path, slides: list[dict[str, Any]]) -> tuple[list[str], list[str], dict[str, Any]]:
    failures: list[str] = []
    warnings: list[str] = []
    evidence: dict[str, Any] = {}
    spec_path = root / "spec_lock.json"
    svg_present = has_svg_pages(root)
    evidence["svg_pages_present"] = svg_present

    if not spec_path.exists():
        if svg_present:
            failures.append("SVG/PPTX-oriented project needs spec_lock.json with layout_execution_contract")
        elif len(slides) > 8:
            warnings.append("decks over 8 slides should include spec_lock.json with layout_execution_contract")
        return failures, warnings, evidence

    try:
        spec = load_json(spec_path)
    except Exception as exc:
        return [f"invalid spec_lock.json: {exc}"], warnings, evidence
    if not isinstance(spec, dict):
        return ["spec_lock.json must be an object"], warnings, evidence

    contract = spec.get("layout_execution_contract")
    if not isinstance(contract, dict):
        if svg_present or len(slides) > 8:
            failures.append("spec_lock.json missing layout_execution_contract")
        else:
            warnings.append("spec_lock.json should include layout_execution_contract before rendering")
        return failures, warnings, evidence

    contract_slides = normalize_contract_slides(contract)
    evidence["layout_contract_slide_count"] = len(contract_slides)
    if slides and len(contract_slides) < len(slides):
        failures.append("layout_execution_contract needs a slide entry for every slide")

    coordinate_policy = str(contract.get("coordinate_policy") or "").lower()
    if coordinate_policy and "absolute" not in coordinate_policy:
        warnings.append("layout_execution_contract.coordinate_policy should use absolute stage coordinates")
    text_fit_policy = str(contract.get("text_fit_policy") or "").lower()
    if not text_fit_policy:
        warnings.append("layout_execution_contract should declare text_fit_policy")
    elif not any(term in text_fit_policy for term in ("line", "fit", "shrink", "split", "overflow", "foreignobject")):
        warnings.append("layout_execution_contract.text_fit_policy should explain line breaks or overflow handling")

    known_layout_ids = layout_pattern_ids()
    for idx, slide in enumerate(contract_slides, start=1):
        slide_label = str(slide.get("slide_id") or slide.get("slide_no") or slide.get("page") or idx)
        missing = [field for field in LAYOUT_CONTRACT_SLIDE_FIELDS if not slide.get(field)]
        if missing:
            failures.append(f"layout_execution_contract slide {slide_label} missing fields: {', '.join(missing)}")

        rhythm = str(slide.get("rhythm") or "").lower()
        if rhythm and rhythm not in ALLOWED_PAGE_RHYTHMS:
            failures.append(
                f"layout_execution_contract slide {slide_label} rhythm must be one of: "
                + ", ".join(sorted(ALLOWED_PAGE_RHYTHMS))
            )

        layout_pattern = str(slide.get("layout_pattern_id") or "")
        declared_layout_ids = {match.group(0).upper() for match in re.finditer(r"\bL\d{2}\b", layout_pattern.upper())}
        if layout_pattern and not declared_layout_ids:
            failures.append(f"layout_execution_contract slide {slide_label} layout_pattern_id should include an Lxx pattern id")
        unknown_layout_ids = declared_layout_ids - known_layout_ids
        if unknown_layout_ids:
            failures.append(
                f"layout_execution_contract slide {slide_label} has unknown layout ids: "
                + ", ".join(sorted(unknown_layout_ids))
            )

        coordinate_slots = slide.get("coordinate_slots")
        if isinstance(coordinate_slots, list) and coordinate_slots:
            slot_ids: list[str] = []
            for slot_idx, slot in enumerate(coordinate_slots, start=1):
                if not isinstance(slot, dict):
                    failures.append(f"layout_execution_contract slide {slide_label} coordinate slot {slot_idx} must be an object")
                    continue
                missing_slot = [
                    field for field in COORDINATE_SLOT_FIELDS if slot.get(field) in ("", None)
                ]
                if missing_slot:
                    failures.append(
                        f"layout_execution_contract slide {slide_label} coordinate slot {slot_idx} missing fields: "
                        + ", ".join(missing_slot)
                    )
                for field in ("x", "y", "w", "h"):
                    if field in slot and slot.get(field) not in ("", None):
                        try:
                            value = float(slot[field])
                        except Exception:
                            failures.append(
                                f"layout_execution_contract slide {slide_label} slot {slot.get('slot_id', slot_idx)} {field} must be numeric"
                            )
                            continue
                        if field in {"w", "h"} and value <= 0:
                            failures.append(
                                f"layout_execution_contract slide {slide_label} slot {slot.get('slot_id', slot_idx)} {field} must be positive"
                            )
                slot_ids.append(str(slot.get("slot_id") or "").strip().lower())
            if slot_ids and not any("title" in slot_id or "claim" in slot_id for slot_id in slot_ids):
                warnings.append(f"layout_execution_contract slide {slide_label} has no title/claim coordinate slot")
            if slot_ids and not any(any(token in slot_id for token in PROOF_SLOT_TOKENS) for slot_id in slot_ids):
                warnings.append(f"layout_execution_contract slide {slide_label} has no obvious proof/media coordinate slot")

        group_ids = slide.get("group_ids")
        declared_groups: set[str] = set()
        if isinstance(group_ids, list):
            declared_groups = {str(group_id).strip() for group_id in group_ids if str(group_id).strip()}
            if len(declared_groups) != len(group_ids):
                warnings.append(f"layout_execution_contract slide {slide_label} has blank or duplicate group_ids")
        elif group_ids:
            failures.append(f"layout_execution_contract slide {slide_label} group_ids must be a list")

        if svg_present and declared_groups:
            candidates = svg_candidates_for_slide(root, slide)
            if not candidates:
                warnings.append(f"layout_execution_contract slide {slide_label} has group_ids but no matching SVG page found")
            else:
                svg_groups: set[str] = set()
                top_counts: list[str] = []
                for candidate in candidates:
                    svg_groups.update(group_ids_in_svg(candidate))
                    top_counts.append(f"{candidate.relative_to(root)}={top_level_group_count(candidate)}")
                missing_groups = sorted(declared_groups - svg_groups)
                if missing_groups:
                    failures.append(
                        f"layout_execution_contract slide {slide_label} declares group_ids missing from SVG: "
                        + ", ".join(missing_groups)
                    )
                low_top_group_files = [item for item in top_counts if item.endswith("=0") or item.endswith("=1")]
                if low_top_group_files:
                    warnings.append(
                        f"layout_execution_contract slide {slide_label} SVG may lack semantic top-level groups: "
                        + ", ".join(low_top_group_files)
                    )
                proof_words = (
                    str(slide.get("proof_object") or "")
                    + " "
                    + str(slide.get("component_type") or "")
                ).lower()
                svg_element_ids: set[str] = set()
                for candidate in candidates:
                    svg_element_ids.update(element_ids_in_svg(candidate))
                if "chart" in proof_words and "chart-plot-area" not in svg_element_ids:
                    warnings.append(
                        f"layout_execution_contract slide {slide_label} is chart-like but SVG lacks chart-plot-area group/id"
                    )

    animations_path = root / "animations.json"
    if animations_path.exists():
        try:
            animations = load_json(animations_path)
        except Exception as exc:
            failures.append(f"invalid animations.json: {exc}")
            animations = None
        if isinstance(animations, dict):
            evidence["animations_json"] = "animations.json"
            for slide_key, groups in animation_targets(animations):
                if not groups:
                    continue
                number_match = re.search(r"\d{1,3}", slide_key)
                pseudo_slide: dict[str, Any] = {"slide_id": slide_key}
                if number_match:
                    pseudo_slide["slide_no"] = int(number_match.group(0))
                candidates = svg_candidates_for_slide(root, pseudo_slide)
                if not candidates:
                    warnings.append(f"animations.json references slide {slide_key} but no matching SVG page was found")
                    continue
                svg_groups = set()
                for candidate in candidates:
                    svg_groups.update(group_ids_in_svg(candidate))
                missing_animation_groups = sorted(set(groups) - svg_groups)
                if missing_animation_groups:
                    failures.append(
                        f"animations.json slide {slide_key} targets missing SVG groups: "
                        + ", ".join(missing_animation_groups)
                    )

    evidence["svg_group_ids"] = {
        rel: sorted(group_ids)
        for rel, group_ids in all_svg_group_ids(root).items()
    }
    return failures, warnings, evidence


def scan_text(path: Path, patterns: list[str]) -> list[str]:
    text = strip_non_visible_payloads(read_text(path))
    hits: list[str] = []
    for pattern in patterns:
        if re.search(pattern, text, flags=re.IGNORECASE):
            hits.append(pattern)
    return hits


def strip_non_visible_payloads(text: str) -> str:
    """Remove encoded media payloads before visible-text policy scans."""
    return DATA_URI_PATTERN.sub("data:payload-removed", text)


def scan_visible_metadata_text(text: str) -> list[str]:
    hits: list[str] = []
    for label, pattern in VISIBLE_INTERNAL_METADATA:
        if re.search(pattern, text, flags=re.IGNORECASE):
            hits.append(label)
    return hits


def visible_policy_patterns(root: Path) -> list[tuple[str, str]]:
    patterns = list(VISIBLE_INTERNAL_METADATA)
    policy_path = root / "visible_provenance_policy.json"
    policy = None
    if not policy_path.exists():
        patterns.extend(VISIBLE_PRODUCTION_JARGON)
        return patterns
    try:
        policy = load_json(policy_path)
    except Exception:
        patterns.extend(VISIBLE_PRODUCTION_JARGON)
        return patterns
    if not isinstance(policy, dict):
        patterns.extend(VISIBLE_PRODUCTION_JARGON)
        return patterns
    if not truthy(policy.get("allow_visible_production_jargon")):
        patterns.extend(VISIBLE_PRODUCTION_JARGON)
    forbidden = policy.get("forbidden_visible_strings")
    if not isinstance(forbidden, list):
        return patterns
    for item in forbidden:
        value = str(item).strip()
        if not value:
            continue
        label = value
        if value.startswith("re:"):
            patterns.append((label, value[3:]))
        else:
            patterns.append((label, re.escape(value)))
    return patterns


def scan_visible_policy_text(text: str, patterns: list[tuple[str, str]]) -> list[str]:
    hits: list[str] = []
    for label, pattern in patterns:
        if re.search(pattern, text, flags=re.IGNORECASE):
            hits.append(label)
    return hits


def html_looks_like_whole_slide_preview(text: str) -> bool:
    return any(re.search(pattern, text, flags=re.IGNORECASE) for pattern in FORMAL_HTML_SCREENSHOT_PATTERNS)


def numeric_line_height(value: str, unit: str) -> float | None:
    try:
        number = float(value)
    except Exception:
        return None
    unit = unit.strip().lower()
    if unit == "%":
        return number / 100
    if unit:
        return None
    return number


def css_block_title_related(selector: str) -> bool:
    normalized = selector.strip().lower()
    return any(token in normalized for token in TITLE_SELECTOR_TOKENS)


def check_html_title_line_heights(path: Path, text: str) -> tuple[list[str], list[str]]:
    failures: list[str] = []
    warnings: list[str] = []
    style_blocks = [
        match.group(1)
        for match in re.finditer(r"<style\b[^>]*>(.*?)</style>", text, flags=re.IGNORECASE | re.DOTALL)
    ]
    if not style_blocks:
        return failures, warnings
    css_text = "\n".join(style_blocks)
    for match in re.finditer(r"([^{}]+)\{([^{}]*)\}", css_text, flags=re.DOTALL):
        selector = match.group(1).strip()
        if not css_block_title_related(selector):
            continue
        body = match.group(2)
        candidates: list[tuple[str, str, str]] = []
        for item in re.finditer(
            r"line-height\s*:\s*([0-9]*\.?[0-9]+)\s*(%|[a-z]+)?",
            body,
            flags=re.IGNORECASE,
        ):
            candidates.append(("line-height", item.group(1), item.group(2) or ""))
        for item in re.finditer(
            r"font\s*:\s*[^;{}]*?/([0-9]*\.?[0-9]+)\s*(%|[a-z]+)?",
            body,
            flags=re.IGNORECASE,
        ):
            candidates.append(("font shorthand line-height", item.group(1), item.group(2) or ""))
        for source, value, unit in candidates:
            normalized = numeric_line_height(value, unit)
            if normalized is None:
                continue
            label = f"{path.relative_to(path.parents[0]) if path.is_absolute() else path} selector `{selector}` {source}={value}{unit}"
            if normalized < MIN_CJK_TITLE_LINE_HEIGHT:
                failures.append(
                    f"title line-height too tight for CJK multi-line titles: {label}; use >= {MIN_CJK_TITLE_LINE_HEIGHT:.2f}, normally {WARN_CJK_TITLE_LINE_HEIGHT:.2f}-1.18"
                )
            elif normalized < WARN_CJK_TITLE_LINE_HEIGHT:
                warnings.append(
                    f"title line-height is tight; verify visually: {label}; normal CJK title range is {WARN_CJK_TITLE_LINE_HEIGHT:.2f}-1.18"
                )
    return failures, warnings


def check_html_delivery_manifest(root: Path, slides: list[dict[str, Any]]) -> tuple[list[str], list[str], dict[str, Any]]:
    failures: list[str] = []
    warnings: list[str] = []
    evidence: dict[str, Any] = {}
    manifest_path = root / "html_delivery_manifest.json"
    if not manifest_path.exists():
        return failures, warnings, evidence
    try:
        manifest = load_json(manifest_path)
    except Exception as exc:
        return [f"invalid html_delivery_manifest.json: {exc}"], warnings, evidence
    if not isinstance(manifest, dict):
        return ["html_delivery_manifest.json must be an object"], warnings, evidence
    evidence["html_delivery_manifest"] = manifest
    if manifest.get("mode") != "semantic_html_deck":
        failures.append("html_delivery_manifest.json mode must be semantic_html_deck")
    try:
        html_count = int(manifest.get("slide_count", 0))
    except Exception:
        html_count = 0
    if slides and html_count != len(slides):
        failures.append(
            f"html_delivery_manifest.json slide_count {html_count} does not match slide_plan count {len(slides)}"
        )
    outputs = manifest.get("html_outputs")
    if not isinstance(outputs, list) or not outputs:
        failures.append("html_delivery_manifest.json needs non-empty html_outputs")
        outputs = []
    component_strategy = str(manifest.get("component_strategy") or "").lower()
    if not any(term in component_strategy for term in ("dom", "svg", "canvas", "css", "js")):
        failures.append("html_delivery_manifest.json component_strategy must mention DOM/SVG/Canvas/CSS/JS")
    screenshot_policy = str(manifest.get("whole_slide_screenshot_policy") or "").lower()
    if not any(term in screenshot_policy for term in ("forbid", "forbidden", "no", "not")):
        failures.append("html_delivery_manifest.json whole_slide_screenshot_policy must forbid whole-slide screenshots")
    motion_system = manifest.get("motion_system")
    motion_manifest_default = root / "html_motion_manifest.json"
    if isinstance(motion_system, dict):
        level = str(motion_system.get("level") or "none").lower()
        manifest_rel = str(motion_system.get("manifest") or "").strip()
        engines = motion_system.get("engines")
        if level not in {"none", "subtle", "expressive", "cinematic"}:
            failures.append("html_delivery_manifest.motion_system.level must be none, subtle, expressive, or cinematic")
        if level != "none":
            if not manifest_rel:
                failures.append("html_delivery_manifest.motion_system.manifest must point to html_motion_manifest.json")
            else:
                motion_manifest_path = root / manifest_rel
                if not motion_manifest_path.exists():
                    failures.append(f"html_delivery_manifest.motion_system manifest missing file: {manifest_rel}")
                else:
                    evidence["html_motion_manifest"] = manifest_rel
            if not isinstance(engines, list) or not engines:
                failures.append("html_delivery_manifest.motion_system.engines must list the authored motion engines")
            fallback = str(motion_system.get("fallback") or "").lower()
            if not any(term in fallback for term in ("static", "fallback", "reduced", "readable")):
                failures.append("html_delivery_manifest.motion_system.fallback must describe readable static/reduced-motion behavior")
    elif motion_manifest_default.exists():
        warnings.append("html_motion_manifest.json exists but html_delivery_manifest.json has no motion_system block")
    readability = manifest.get("readability_qa")
    if not isinstance(readability, dict):
        failures.append("html_delivery_manifest.json needs readability_qa for formal HTML")
        readability = {}
    else:
        viewports = readability.get("viewports_checked")
        if not isinstance(viewports, list) or len(viewports) < 2:
            failures.append("html_delivery_manifest.readability_qa.viewports_checked needs at least two viewport sizes")
        stage_strategy = str(readability.get("stage_strategy") or "").lower()
        if not any(term in stage_strategy for term in ("16:9", "aspect", "scale", "1920", "1280")):
            failures.append("html_delivery_manifest.readability_qa.stage_strategy must describe fixed 16:9 responsive scaling")
        if "1920" in stage_strategy and not any(
            term in stage_strategy for term in ("left", "top", "origin", "explicit", "position")
        ):
            failures.append(
                "html_delivery_manifest.readability_qa.stage_strategy must describe explicit positioning when using a 1920px coordinate stage"
            )
        decoration_budget = readability.get("background_decoration_budget")
        if not isinstance(decoration_budget, dict):
            warnings.append(
                "html_delivery_manifest.readability_qa.background_decoration_budget should record line/noise caps for formal HTML backgrounds"
            )
        else:
            level = str(decoration_budget.get("level") or "").lower()
            if level and level not in {"quiet", "moderate", "cinematic"}:
                failures.append(
                    "html_delivery_manifest.readability_qa.background_decoration_budget.level must be quiet, moderate, or cinematic"
                )
            for field in (
                "decorative_line_families_max_body",
                "standalone_line_segments_max_body",
                "safe_area_clearance_px_min_1280x720",
            ):
                if field not in decoration_budget:
                    warnings.append(
                        f"html_delivery_manifest.readability_qa.background_decoration_budget missing {field}"
                    )
        try:
            min_body = int(readability.get("min_body_px_at_1280_stage", 0))
        except Exception:
            min_body = 0
        if min_body < 18:
            failures.append("html_delivery_manifest.readability_qa.min_body_px_at_1280_stage must be >= 18")
        overflow_policy = str(readability.get("overflow_policy") or "").lower()
        if not any(term in overflow_policy for term in ("no hidden", "no clipping", "visible", "fit", "scroll")):
            failures.append("html_delivery_manifest.readability_qa.overflow_policy must prevent hidden/clipped slide content")
        parity_policy = str(readability.get("content_parity_policy") or "").lower()
        if not any(term in parity_policy for term in ("slide_plan", "same", "parity", "titles", "anchors")):
            failures.append("html_delivery_manifest.readability_qa.content_parity_policy must require slide-plan content parity")
        browser_screenshots = readability.get("browser_screenshots")
        if not isinstance(browser_screenshots, list) or not browser_screenshots:
            failures.append("html_delivery_manifest.readability_qa.browser_screenshots must list checked browser screenshots")
            browser_screenshots = []
        required_screenshot_count = min(len(slides), 4) if slides else 1
        if len(browser_screenshots) < required_screenshot_count:
            failures.append(
                f"html_delivery_manifest.readability_qa.browser_screenshots needs at least {required_screenshot_count} representative slide screenshots"
            )
        for rel in browser_screenshots:
            screenshot_path = root / str(rel)
            if not screenshot_path.exists():
                failures.append(f"html_delivery_manifest.readability_qa.browser_screenshots missing file: {rel}")
    for rel in outputs:
        path = root / str(rel)
        if not path.exists():
            failures.append(f"html_delivery_manifest.json missing html output: {rel}")
            continue
        if path.suffix.lower() != ".html":
            failures.append(f"html_delivery_manifest output is not HTML: {rel}")
            continue
        text = read_text(path)
        if html_looks_like_whole_slide_preview(text):
            failures.append(
                f"{rel} appears to be whole-slide preview-image HTML; formal HTML must be semantic DOM/SVG/Canvas/CSS/JS"
            )
        missing_titles = [
            slide_title(slide)
            for slide in slides
            if slide_title(slide) and slide_title(slide) not in text
        ]
        if missing_titles:
            failures.append(
                f"{rel} missing slide-plan titles in visible/semantic HTML: "
                + "; ".join(missing_titles[:4])
            )
    return failures, warnings, evidence


def check_preview_gate(root: Path, slides: list[dict[str, Any]]) -> tuple[list[str], list[str], dict[str, Any]]:
    failures: list[str] = []
    warnings: list[str] = []
    evidence: dict[str, Any] = {}
    if len(slides) <= 7:
        return failures, warnings, evidence

    gate_path = root / "preview_gate.json"
    evidence["preview_gate"] = str(gate_path.relative_to(root)) if gate_path.exists() else ""
    if not gate_path.exists():
        failures.append("decks over 7 slides need preview_gate.json with approved four-slide preview or explicit user skip")
        return failures, warnings, evidence
    try:
        gate = load_json(gate_path)
    except Exception as exc:
        return [f"invalid preview_gate.json: {exc}"], warnings, evidence
    if not isinstance(gate, dict):
        return ["preview_gate.json must be an object"], warnings, evidence
    evidence["preview_gate"] = gate
    skipped = gate.get("skipped_by_user") is True
    decision = str(gate.get("user_decision") or "").lower()
    if not skipped and decision != "approved":
        failures.append("preview_gate.json user_decision must be approved before full generation")
    if skipped and not gate.get("skip_instruction"):
        failures.append("preview_gate.json skipped_by_user=true requires skip_instruction")
    selected = gate.get("selected_slides")
    if not isinstance(selected, list) or len(selected) != 4:
        failures.append("preview_gate.json selected_slides must list exactly four slides")
    outputs = gate.get("outputs")
    if not isinstance(outputs, list) or not outputs:
        failures.append("preview_gate.json needs non-empty outputs")
    else:
        missing_outputs = [str(rel) for rel in outputs if not (root / str(rel)).exists()]
        if missing_outputs:
            failures.append("preview_gate.json outputs missing: " + ", ".join(missing_outputs[:6]))
    qa_focus = gate.get("qa_focus")
    if isinstance(qa_focus, list):
        normalized_focus = " ".join(str(item).lower() for item in qa_focus)
        for term in ("typography", "background", "connector", "html"):
            if term not in normalized_focus:
                warnings.append(f"preview_gate.json qa_focus should include {term}")
    else:
        warnings.append("preview_gate.json should list qa_focus items")
    return failures, warnings, evidence


def check_export_manifest(root: Path) -> tuple[list[str], list[str], dict[str, Any]]:
    failures: list[str] = []
    warnings: list[str] = []
    evidence: dict[str, Any] = {}
    manifest_path = root / "export_manifest.json"
    if not manifest_path.exists():
        exports = root / "exports"
        if exports.exists():
            has_multiple_formats = sum(
                1
                for pattern in ("*.pptx", "*.pdf", "*.html", "*.key")
                if list(exports.glob(pattern))
            ) >= 2
            if has_multiple_formats:
                warnings.append("multi-format exports exist but export_manifest.json is missing")
        return failures, warnings, evidence

    try:
        manifest = load_json(manifest_path)
    except Exception as exc:
        return [f"invalid export_manifest.json: {exc}"], warnings, evidence
    if not isinstance(manifest, dict):
        return ["export_manifest.json must be an object"], warnings, evidence

    evidence["export_manifest"] = manifest
    keynote_capability = manifest.get("keynote_capability")
    if isinstance(keynote_capability, dict):
        report_rel = str(keynote_capability.get("report") or "")
        evidence["keynote_capability"] = keynote_capability
        if report_rel and not (root / report_rel).exists():
            failures.append(f"export_manifest.json keynote_capability report missing: {report_rel}")
    formats = manifest.get("formats")
    if not isinstance(formats, dict) or not formats:
        failures.append("export_manifest.json needs non-empty formats object")
        return failures, warnings, evidence

    pptx_path: Path | None = None
    pptx_item = formats.get("pptx")
    if isinstance(pptx_item, dict) and pptx_item.get("path"):
        candidate = root / str(pptx_item.get("path"))
        if candidate.exists():
            pptx_path = candidate

    requested = manifest.get("requested_formats", [])
    if requested and not isinstance(requested, list):
        failures.append("export_manifest.json requested_formats must be a list")
    last_requested = manifest.get("last_requested_formats", [])
    if last_requested and not isinstance(last_requested, list):
        failures.append("export_manifest.json last_requested_formats must be a list")

    allowed_statuses = {"exported", "existing", "missing", "failed"}
    for name, item in formats.items():
        if not isinstance(item, dict):
            failures.append(f"export_manifest.json format {name} must be an object")
            continue
        status = item.get("status")
        if status not in allowed_statuses:
            failures.append(f"export_manifest.json format {name} has invalid status: {status}")
            continue
        if status in {"exported", "existing"}:
            path_value = item.get("path") or item.get("index")
            if not path_value:
                failures.append(f"export_manifest.json successful format {name} needs path or index")
                continue
            rel_path = str(path_value)
            if not (root / rel_path).exists():
                failures.append(f"export_manifest.json format {name} missing output file: {rel_path}")
            output_path = root / rel_path
            if pptx_path and name in {"pdf", "html-parity", "keynote"} and output_path.exists():
                try:
                    if artifact_mtime(output_path) < artifact_mtime(pptx_path):
                        failures.append(
                            f"export_manifest.json format {name} output is older than PPTX source: {rel_path}"
                        )
                except OSError:
                    pass
            if name == "html":
                if rel_path.endswith(".parity.html") or rel_path.startswith("html-parity/"):
                    failures.append("export_manifest.json formal html path must not point to parity HTML")
                if not (root / "html_delivery_manifest.json").exists():
                    failures.append("export_manifest.json reports formal HTML success but html_delivery_manifest.json is missing")
            if name == "html-parity":
                if not (rel_path.endswith(".parity.html") or rel_path.startswith("html-parity/")):
                    failures.append("export_manifest.json html-parity path must use html-parity/ or .parity.html")
                if not (root / "html_parity_manifest.json").exists():
                    failures.append("export_manifest.json reports html-parity success but html_parity_manifest.json is missing")
            if name == "pdf" and not rel_path.endswith(".pdf"):
                failures.append("export_manifest.json pdf output must end with .pdf")
            if name == "keynote" and not rel_path.endswith(".key"):
                failures.append("export_manifest.json keynote output must end with .key")
            if name == "keynote" and item.get("compatibility_format") == "Keynote 09":
                if item.get("fallback_from") != "save as Keynote":
                    failures.append("export_manifest.json Keynote 09 fallback must record fallback_from='save as Keynote'")
                if not isinstance(item.get("primary_failure"), dict) or not item.get("primary_failure"):
                    failures.append("export_manifest.json Keynote 09 fallback must record primary_failure")
        if status in {"missing", "failed"} and not (item.get("reason") or item.get("warning")):
            warnings.append(f"export_manifest.json format {name} is {status} without reason")
        if name == "keynote" and status == "failed" and not item.get("diagnostic_command"):
            warnings.append("export_manifest.json keynote failure should include diagnostic_command")

    return failures, warnings, evidence


def iter_pptx_visible_text(path: Path) -> tuple[list[tuple[int, str]], str | None]:
    try:
        from pptx import Presentation as PptxPresentation
    except Exception as exc:
        return [], f"python-pptx unavailable; cannot inspect visible text in {path.name}: {exc}"

    try:
        presentation = PptxPresentation(str(path))
    except Exception as exc:
        return [], f"cannot inspect visible text in {path.name}: {exc}"

    visible_text: list[tuple[int, str]] = []
    for slide_idx, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text:
                visible_text.append((slide_idx, text))
    return visible_text, None


def inspect_pptx_editability(path: Path) -> tuple[dict[str, Any], str | None]:
    try:
        from pptx import Presentation as PptxPresentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except Exception as exc:
        return {}, f"python-pptx unavailable; cannot inspect editability in {path.name}: {exc}"

    try:
        presentation = PptxPresentation(str(path))
    except Exception as exc:
        return {}, f"cannot inspect editability in {path.name}: {exc}"

    slide_w = int(presentation.slide_width)
    slide_h = int(presentation.slide_height)
    slides: list[dict[str, Any]] = []
    image_backed_count = 0
    for slide_idx, slide in enumerate(presentation.slides, start=1):
        full_slide_pictures = 0
        native_text_chars = 0
        native_text_shapes = 0
        picture_count = 0
        for shape in slide.shapes:
            text = getattr(shape, "text", "") or ""
            if text.strip():
                native_text_shapes += 1
                native_text_chars += len(text.strip())
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                picture_count += 1
                left = int(getattr(shape, "left", 0))
                top = int(getattr(shape, "top", 0))
                width = int(getattr(shape, "width", 0))
                height = int(getattr(shape, "height", 0))
                if width >= slide_w * 0.88 and height >= slide_h * 0.88 and left <= slide_w * 0.06 and top <= slide_h * 0.06:
                    full_slide_pictures += 1
        image_backed = full_slide_pictures > 0 and native_text_chars < 12
        if image_backed:
            image_backed_count += 1
        slides.append(
            {
                "slide_no": slide_idx,
                "picture_count": picture_count,
                "full_slide_picture_count": full_slide_pictures,
                "native_text_shape_count": native_text_shapes,
                "native_text_chars": native_text_chars,
                "image_backed": image_backed,
            }
        )
    slide_count = len(slides)
    return {
        "slide_count": slide_count,
        "image_backed_slide_count": image_backed_count,
        "image_backed_ratio": image_backed_count / slide_count if slide_count else 0,
        "slides": slides,
    }, None


def inspect_pptx_picture_aspects(path: Path) -> tuple[dict[str, Any], str | None]:
    try:
        from pptx import Presentation as PptxPresentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except Exception as exc:
        return {}, f"python-pptx unavailable; cannot inspect picture aspect ratios in {path.name}: {exc}"

    try:
        presentation = PptxPresentation(str(path))
    except Exception as exc:
        return {}, f"cannot inspect picture aspect ratios in {path.name}: {exc}"

    pictures: list[dict[str, Any]] = []
    distorted: list[dict[str, Any]] = []
    for slide_idx, slide in enumerate(presentation.slides, start=1):
        picture_idx = 0
        for shape in slide.shapes:
            if getattr(shape, "shape_type", None) != MSO_SHAPE_TYPE.PICTURE:
                continue
            picture_idx += 1
            width = int(getattr(shape, "width", 0))
            height = int(getattr(shape, "height", 0))
            if width <= 0 or height <= 0:
                continue
            try:
                image_width, image_height = shape.image.size
            except Exception:
                continue
            if image_width <= 0 or image_height <= 0:
                continue
            shape_ratio = width / height
            image_ratio = image_width / image_height
            distortion = abs(shape_ratio - image_ratio) / image_ratio
            record = {
                "slide_no": slide_idx,
                "picture_no": picture_idx,
                "shape_ratio": round(shape_ratio, 4),
                "image_ratio": round(image_ratio, 4),
                "distortion_ratio": round(distortion, 4),
                "shape_inches": {
                    "w": round(width / 914400, 3),
                    "h": round(height / 914400, 3),
                },
                "image_px": {"w": image_width, "h": image_height},
            }
            pictures.append(record)
            if distortion > DEFAULT_MAX_PICTURE_ASPECT_DISTORTION:
                distorted.append(record)
    return {
        "max_allowed_distortion_ratio": DEFAULT_MAX_PICTURE_ASPECT_DISTORTION,
        "picture_count": len(pictures),
        "distorted_count": len(distorted),
        "distorted": distorted,
        "pictures": pictures,
    }, None


def detect_quality_intent(root: Path, slide_count: int) -> dict[str, Any]:
    evidence: dict[str, Any] = {"mode": "normal", "threshold": 0, "signals": []}
    if slide_count <= 7:
        return evidence
    probes = [
        root / "slide_plan.json",
        root / "visual_asset_manifest.json",
        root / "preview_gate.json",
        root / "production_manifest.json",
        root / "route_card.md",
        root / "generation_report.md",
    ]
    blob_parts: list[str] = []
    for path in probes:
        if not path.exists() or not path.is_file():
            continue
        try:
            blob_parts.append(read_text(path)[:8000].lower())
        except Exception:
            continue
    blob = "\n".join(blob_parts)
    if "ppt_master_grade" in blob or "ppt-master" in blob or "ppt_master_ready" in blob:
        evidence["mode"] = "ppt_master_grade"
        evidence["threshold"] = 85
        evidence["signals"].append("ppt_master_grade")
    elif "quality_profile" in blob and "final" in blob:
        evidence["mode"] = "final"
        evidence["threshold"] = 85
        evidence["signals"].append("quality_profile=final")
    elif "quality_profile" in blob and "professional" in blob:
        evidence["mode"] = "professional"
        evidence["threshold"] = 75
        evidence["signals"].append("quality_profile=professional")
    return evidence


def check_quality_benchmark_gate(root: Path, slide_count: int) -> tuple[list[str], list[str], dict[str, Any]]:
    failures: list[str] = []
    warnings: list[str] = []
    intent = detect_quality_intent(root, slide_count)
    benchmark_path = root / "reports" / "deck_quality_benchmark.json"
    repair_path = root / "reports" / "deck_repair_plan.json"
    evidence: dict[str, Any] = {
        "intent": intent,
        "benchmark": str(benchmark_path.relative_to(root)) if benchmark_path.exists() else "",
        "repair_plan": str(repair_path.relative_to(root)) if repair_path.exists() else "",
    }
    if slide_count <= 7:
        return failures, warnings, evidence
    threshold = int(intent.get("threshold") or 0)
    high_intent = threshold > 0
    if not benchmark_path.exists():
        message = "reports/deck_quality_benchmark.json missing; cannot prove content depth, image density, style execution, and layout rhythm"
        if high_intent:
            failures.append(message)
        else:
            warnings.append(message)
        return failures, warnings, evidence
    try:
        benchmark = load_json(benchmark_path)
    except Exception as exc:
        return [f"invalid reports/deck_quality_benchmark.json: {exc}"], warnings, evidence
    if not isinstance(benchmark, dict):
        return ["reports/deck_quality_benchmark.json must be an object"], warnings, evidence
    if benchmark.get("schema_version") != "1.0.0" or not isinstance(benchmark.get("categories"), list) or not isinstance(benchmark.get("stats"), dict):
        return [
            "reports/deck_quality_benchmark.json is not an official qiaomu-ppt deck_quality_benchmark.py report; rerun the skill script instead of hand-writing a benchmark sidecar"
        ], warnings, evidence
    score = int(benchmark.get("score") or 0)
    benchmark_target = int(benchmark.get("target_score") or 0)
    readiness = str(benchmark.get("readiness") or "")
    ppt_master_ready = bool(benchmark.get("ppt_master_ready"))
    evidence.update(
        {
            "score": score,
            "target_score": benchmark_target,
            "readiness": readiness,
            "ppt_master_ready": ppt_master_ready,
        }
    )
    effective_threshold = max(threshold, benchmark_target if benchmark_target >= 75 else 0)
    benchmark_enforced = high_intent or benchmark_target >= 75
    if benchmark_enforced and score < effective_threshold:
        failures.append(f"deck_quality_benchmark score {score} below required {effective_threshold} for {intent['mode']}")
    elif not high_intent and score < 70:
        warnings.append(f"deck_quality_benchmark score {score} is weak; review repair plan before final delivery")
    if intent.get("mode") == "ppt_master_grade" and not ppt_master_ready:
        failures.append("project claims ppt_master_grade but benchmark does not mark ppt_master_ready")

    if repair_path.exists():
        try:
            repair = load_json(repair_path)
        except Exception as exc:
            failures.append(f"invalid reports/deck_repair_plan.json: {exc}")
            return failures, warnings, evidence
        actions = repair.get("actions") if isinstance(repair, dict) else []
        critical = [
            item
            for item in actions
            if isinstance(item, dict) and str(item.get("severity") or "").lower() == "critical"
        ] if isinstance(actions, list) else []
        evidence["critical_repair_count"] = len(critical)
        if benchmark_enforced and critical:
            failures.append(
                "deck_repair_plan contains critical repair actions: "
                + ", ".join(str(item.get("action_id") or item.get("title") or "critical") for item in critical[:5])
            )
    elif high_intent:
        warnings.append("reports/deck_repair_plan.json missing; high-design/final decks should record repair actions after benchmark")
    return failures, warnings, evidence


def check_project(root: Path, *, require_real_imagegen: bool = False) -> dict[str, Any]:
    root = root.resolve()
    failures: list[str] = []
    warnings: list[str] = []
    evidence: dict[str, Any] = {"root": str(root)}
    visible_text_patterns = visible_policy_patterns(root)

    for name in REQUIRED_ROOT:
        if not (root / name).exists():
            failures.append(f"missing required artifact: {name}")

    deck_brief_text = read_text(root / "deck_brief.md") if (root / "deck_brief.md").exists() else ""
    source_failures, source_warnings, source_evidence = check_source_manifest(root, deck_brief_text)
    failures.extend(source_failures)
    warnings.extend(source_warnings)
    evidence["source_manifest"] = source_evidence

    recommendations_path = root / "style_recommendations.json"
    if recommendations_path.exists():
        try:
            recommendations = load_json(recommendations_path)
            top = recommendations.get("top", []) if isinstance(recommendations, dict) else []
            if not top:
                warnings.append("style_recommendations.json exists but has no top recommendations")
            style_brief = root / "style_brief.md"
            spec_lock = root / "spec_lock.json"
            if style_brief.exists() and "selected_preset" not in read_text(style_brief):
                warnings.append("style recommendations exist but style_brief.md does not mention selected_preset")
            if spec_lock.exists() and "selected_preset" not in read_text(spec_lock):
                warnings.append("style recommendations exist but spec_lock.json does not mention selected_preset")
        except Exception as exc:
            failures.append(f"invalid style_recommendations.json: {exc}")

    plan_path = root / "slide_plan.json"
    slides: list[dict[str, Any]] = []
    if plan_path.exists():
        try:
            slides = iter_slides(load_json(plan_path))
        except Exception as exc:
            failures.append(f"invalid slide_plan.json: {exc}")
        if not slides:
            failures.append("slide_plan.json has no slides")
        for idx, slide in enumerate(slides, start=1):
            missing = [field for field in SLIDE_FIELDS if not slide.get(field)]
            if missing:
                failures.append(f"slide {idx} missing fields: {', '.join(missing)}")
    evidence["slide_count"] = len(slides)

    visual_contract_path = root / "visual_contract.json"
    if len(slides) > 8 and not visual_contract_path.exists():
        failures.append("missing visual_contract.json for deck with more than 8 slides")
    if visual_contract_path.exists():
        contract_failures, contract_warnings, contract_evidence = check_visual_contract(
            visual_contract_path, len(slides)
        )
        failures.extend(contract_failures)
        warnings.extend(contract_warnings)
        evidence["visual_contract"] = contract_evidence

    export_artifacts_exist = (root / "exports").exists() and any((root / "exports").iterdir())
    if len(slides) > 8 and not (root / "visual_asset_manifest.json").exists():
        failures.append("missing visual_asset_manifest.json for deck with more than 8 slides")
    asset_failures, asset_warnings, asset_evidence = check_visual_asset_manifest(
        root, require_terminal=export_artifacts_exist, require_real_imagegen=require_real_imagegen
    )
    failures.extend(asset_failures)
    warnings.extend(asset_warnings)
    if asset_evidence:
        evidence["visual_assets"] = asset_evidence
    source_image_warnings, source_image_evidence = check_source_image_usage(root, asset_evidence)
    warnings.extend(source_image_warnings)
    if source_image_evidence:
        evidence["source_image_usage"] = source_image_evidence

    content_contract_path = root / "content_contract.json"
    content_contract_payload: dict[str, Any] | None = None
    if len(slides) > 8 and not content_contract_path.exists():
        failures.append("missing content_contract.json for deck with more than 8 slides")
    if content_contract_path.exists():
        try:
            loaded_content_contract = load_json(content_contract_path)
            if isinstance(loaded_content_contract, dict):
                content_contract_payload = loaded_content_contract
        except Exception:
            content_contract_payload = None
        content_failures, content_warnings, content_evidence = check_content_contract(
            content_contract_path, slides
        )
        failures.extend(content_failures)
        warnings.extend(content_warnings)
        evidence["content_contract"] = content_evidence

    topic_failures, topic_warnings, topic_evidence = check_topic_research_artifacts(
        root, slides, content_contract_payload
    )
    failures.extend(topic_failures)
    warnings.extend(topic_warnings)
    if topic_evidence:
        evidence["topic_research"] = topic_evidence

    layout_failures, layout_warnings, layout_evidence = check_layout_execution_contract(root, slides)
    failures.extend(layout_failures)
    warnings.extend(layout_warnings)
    evidence["layout_execution"] = layout_evidence

    html_delivery_failures, html_delivery_warnings, html_delivery_evidence = check_html_delivery_manifest(
        root, slides
    )
    failures.extend(html_delivery_failures)
    warnings.extend(html_delivery_warnings)
    evidence.update(html_delivery_evidence)

    preview_failures, preview_warnings, preview_evidence = check_preview_gate(root, slides)
    failures.extend(preview_failures)
    warnings.extend(preview_warnings)
    evidence.update(preview_evidence)

    export_manifest_failures, export_manifest_warnings, export_manifest_evidence = check_export_manifest(root)
    failures.extend(export_manifest_failures)
    warnings.extend(export_manifest_warnings)
    evidence.update(export_manifest_evidence)

    benchmark_failures, benchmark_warnings, benchmark_evidence = check_quality_benchmark_gate(root, len(slides))
    failures.extend(benchmark_failures)
    warnings.extend(benchmark_warnings)
    evidence["quality_benchmark_gate"] = benchmark_evidence

    for path in sorted(root.rglob("*")):
        if not path.is_file():
            continue
        if path.suffix.lower() in {".md", ".json", ".html", ".svg"}:
            placeholder_hits = scan_text(path, PLACEHOLDERS)
            if placeholder_hits:
                failures.append(f"{path.relative_to(root)} has placeholder residue: {', '.join(placeholder_hits)}")
        if path.suffix.lower() == ".svg":
            banned_hits = scan_text(path, SVG_BANNED)
            if banned_hits:
                failures.append(f"{path.relative_to(root)} has SVG/PPTX risky features: {', '.join(banned_hits)}")
            text = read_text(path)
            if "viewBox" not in text:
                warnings.append(f"{path.relative_to(root)} missing viewBox")
        if path.suffix.lower() == ".html":
            text = read_text(path)
            if "1920" not in text or "1080" not in text:
                warnings.append(f"{path.relative_to(root)} may not declare a fixed 1920x1080 stage")
            rel_path = path.relative_to(root).as_posix()
            if not rel_path.startswith("html-parity/") and not path.name.endswith(".parity.html"):
                if html_looks_like_whole_slide_preview(text):
                    failures.append(
                        f"{path.relative_to(root)} appears to be whole-slide preview-image HTML; use html-parity/ or .parity.html, or generate a formal semantic HTML deck"
                    )
            metadata_hits = scan_visible_policy_text(text, visible_text_patterns)
            if metadata_hits:
                failures.append(
                    f"{path.relative_to(root)} contains forbidden visible production text: {', '.join(metadata_hits)}"
                )
            if not rel_path.startswith("html-parity/") and not path.name.endswith(".parity.html"):
                title_line_failures, title_line_warnings = check_html_title_line_heights(path, text)
                failures.extend(
                    f"{path.relative_to(root)}: {failure}" for failure in title_line_failures
                )
                warnings.extend(
                    f"{path.relative_to(root)}: {warning}" for warning in title_line_warnings
                )
        if path.suffix.lower() == ".svg":
            metadata_hits = scan_visible_policy_text(read_text(path), visible_text_patterns)
            if metadata_hits:
                failures.append(
                    f"{path.relative_to(root)} contains forbidden visible production text: {', '.join(metadata_hits)}"
                )

    exports = list((root / "exports").glob("*.pptx")) if (root / "exports").exists() else []
    evidence["pptx_exports"] = [str(path.relative_to(root)) for path in exports]
    html_exports = []
    if (root / "exports").exists():
        html_exports.extend(
            path for path in (root / "exports").glob("*.html") if not path.name.endswith(".parity.html")
        )
    if (root / "html").exists():
        html_exports.extend((root / "html").glob("*.html"))
    parity_html_exports = []
    if (root / "exports").exists():
        parity_html_exports.extend((root / "exports").glob("*.parity.html"))
    if (root / "html-parity").exists():
        parity_html_exports.extend((root / "html-parity").glob("*.html"))
    evidence["html_exports"] = [str(path.relative_to(root)) for path in html_exports]
    evidence["html_parity_exports"] = [str(path.relative_to(root)) for path in parity_html_exports]
    if not exports:
        warnings.append("no PPTX export found; mark editable PPTX export as missing evidence if required")
    if len(slides) > 0 and not html_exports:
        warnings.append("no formal semantic HTML presentation export found; normal PPT runs should produce one unless explicitly forbidden")

    parity_manifest = root / "html_parity_manifest.json"
    if parity_manifest.exists():
        try:
            parity = load_json(parity_manifest)
            evidence["html_parity_manifest"] = parity
            if parity.get("mode") != "rendered_slide_parity":
                failures.append("html_parity_manifest.json mode must be rendered_slide_parity")
            try:
                parity_count = int(parity.get("slide_count", 0))
            except Exception:
                parity_count = 0
            if slides and parity_count != len(slides):
                failures.append(
                    f"html_parity_manifest.json slide_count {parity_count} does not match slide_plan count {len(slides)}"
                )
            for rel in parity.get("html_outputs", []):
                if not (root / str(rel)).exists():
                    failures.append(f"html_parity_manifest.json missing html output: {rel}")
                rel_text = str(rel)
                if rel_text == "html/index.html" or (
                    rel_text.startswith("exports/") and rel_text.endswith(".html") and not rel_text.endswith(".parity.html")
                ):
                    failures.append(
                        "html_parity_manifest.json outputs must use html-parity/ or .parity.html paths, not formal HTML deck paths"
                    )
            for rel in parity.get("preview_images", []):
                if not (root / str(rel)).exists():
                    failures.append(f"html_parity_manifest.json missing preview image: {rel}")
        except Exception as exc:
            failures.append(f"invalid html_parity_manifest.json: {exc}")
    elif exports and html_exports and (root / "previews").exists():
        warnings.append("PPTX-first project has previews and HTML but no html_parity_manifest.json; verify HTML/PPTX content parity manually")
    if exports:
        preview_candidates = list((root / "previews").rglob("*.jpg")) + list((root / "previews").rglob("*.png")) + list((root / "previews").rglob("*.pdf"))
        evidence["preview_artifacts"] = [str(path.relative_to(root)) for path in preview_candidates[:20]]
        if not preview_candidates:
            warnings.append("PPTX export exists but no preview render found under previews/")
        text_check_path = root / "pptx_text_check.json"
        if text_check_path.exists():
            try:
                text_check = load_json(text_check_path)
                evidence["pptx_text_check"] = {
                    "ok": text_check.get("ok"),
                    "path": str(text_check_path.relative_to(root)),
                    "editability": text_check.get("editability", {}),
                }
                if text_check.get("ok") is False:
                    failures.append("pptx_text_check.json reports failures")
                editability = text_check.get("editability", {})
                if isinstance(editability, dict) and editability.get("image_backed_ratio", 0) >= 0.8:
                    failures.append("PPTX text check reports a mostly image-backed deck; normal editable PPTX must use native foreground objects")
            except Exception as exc:
                failures.append(f"invalid pptx_text_check.json: {exc}")
        else:
            failures.append("PPTX export exists but pptx_text_check.json is missing")
        visible_metadata_hits: list[str] = []
        editability_reports: dict[str, Any] = {}
        picture_aspect_reports: dict[str, Any] = {}
        for export in exports:
            editability, editability_warning = inspect_pptx_editability(export)
            if editability_warning:
                warnings.append(editability_warning)
            else:
                editability_reports[str(export.relative_to(root))] = editability
                if editability.get("image_backed_ratio", 0) >= 0.8:
                    failures.append(
                        f"{export.relative_to(root)} is mostly whole-slide raster images; label it as image-backed preview/social output, not editable PPTX"
                    )
            picture_aspects, picture_aspect_warning = inspect_pptx_picture_aspects(export)
            if picture_aspect_warning:
                warnings.append(picture_aspect_warning)
            else:
                picture_aspect_reports[str(export.relative_to(root))] = picture_aspects
                distorted = picture_aspects.get("distorted")
                if isinstance(distorted, list) and distorted:
                    examples = "; ".join(
                        f"slide {item.get('slide_no')} picture {item.get('picture_no')} "
                        f"shape={item.get('shape_ratio')} image={item.get('image_ratio')}"
                        for item in distorted[:6]
                    )
                    failures.append(
                        "PPTX contains stretched picture(s); crop/contain to the slot aspect ratio instead of resizing non-proportionally: "
                        + examples
                    )
            visible_text, inspection_warning = iter_pptx_visible_text(export)
            if inspection_warning:
                warnings.append(inspection_warning)
                continue
            for slide_no, text in visible_text:
                metadata_hits = scan_visible_policy_text(text, visible_text_patterns)
                if metadata_hits:
                    visible_metadata_hits.append(
                        f"{export.relative_to(root)} slide {slide_no}: {', '.join(metadata_hits)}"
                    )
        if editability_reports:
            evidence["pptx_editability"] = editability_reports
        if picture_aspect_reports:
            evidence["pptx_picture_aspects"] = picture_aspect_reports
        evidence["visible_metadata_hits"] = visible_metadata_hits
        if visible_metadata_hits:
            failures.append(
                "PPTX contains forbidden visible production text; move workflow metadata/jargon to reports or notes"
            )

    return {
        "ok": not failures,
        "failures": failures,
        "warnings": warnings,
        "evidence": evidence,
    }


def main() -> None:
    parser = argparse.ArgumentParser(description="Check qiaomu-ppt project artifacts.")
    parser.add_argument("project_dir", help="Generated qiaomu-ppt project directory.")
    parser.add_argument("--output", "-o", help="Write JSON report to this path.")
    parser.add_argument(
        "--require-real-imagegen",
        action="store_true",
        help="Fail if exported AI assets still use procedural-preview-fallback instead of a real image generator.",
    )
    args = parser.parse_args()

    result = check_project(Path(args.project_dir), require_real_imagegen=args.require_real_imagegen)
    rendered = json.dumps(result, ensure_ascii=False, indent=2)
    print(rendered)
    if args.output:
        output = Path(args.output)
        output.parent.mkdir(parents=True, exist_ok=True)
        output.write_text(rendered + "\n", encoding="utf-8")
    if not result["ok"]:
        sys.exit(2)


if __name__ == "__main__":
    main()

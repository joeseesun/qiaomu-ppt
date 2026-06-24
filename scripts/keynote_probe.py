#!/usr/bin/env python3
"""Diagnose macOS Keynote import/save behavior for a generated PPTX.

This script is deliberately separate from export_bundle.py. The exporter should
stay conservative and quick; the probe can spend extra time identifying whether
Keynote automation is unavailable, slow to import, unable to save, or leaving
documents open after an AppleScript timeout.
"""

from __future__ import annotations

import argparse
import json
import platform
import shutil
import subprocess
import time
from datetime import datetime, timezone
from pathlib import Path
from typing import Any


STATUS_OK = "ok"
STATUS_SKIPPED = "skipped"
STATUS_FAILED = "failed"
STATUS_TIMEOUT = "timeout"


def write_json(path: Path, payload: Any) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(payload, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def run_osascript(script: str, timeout: int) -> dict[str, Any]:
    started = time.monotonic()
    try:
        proc = subprocess.run(
            ["osascript", "-e", script],
            text=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            check=False,
            timeout=timeout,
        )
        status = STATUS_OK if proc.returncode == 0 else STATUS_FAILED
        return {
            "status": status,
            "returncode": proc.returncode,
            "duration_sec": round(time.monotonic() - started, 3),
            "stdout": proc.stdout.strip(),
            "stderr": proc.stderr.strip(),
        }
    except subprocess.TimeoutExpired as exc:
        return {
            "status": STATUS_TIMEOUT,
            "duration_sec": round(time.monotonic() - started, 3),
            "stdout": (exc.stdout or "").strip() if isinstance(exc.stdout, str) else "",
            "stderr": (exc.stderr or "").strip() if isinstance(exc.stderr, str) else "",
        }


def applescript_string(value: str) -> str:
    return '"' + value.replace("\\", "\\\\").replace('"', '\\"') + '"'


def applescript_path(path: Path) -> str:
    return str(path).replace("\\", "\\\\").replace('"', '\\"')


def applescript_list(values: list[str]) -> str:
    return "{" + ", ".join(applescript_string(value) for value in values) + "}"


def keynote_document_names(timeout: int = 15) -> list[str]:
    if platform.system() != "Darwin" or not shutil.which("osascript"):
        return []
    result = run_osascript('tell application "Keynote" to get name of documents', timeout)
    if result.get("status") != STATUS_OK:
        return []
    text = str(result.get("stdout") or "").strip()
    if not text:
        return []
    return [item.strip() for item in text.split(",") if item.strip()]


def cleanup_documents(name_prefix: str, before_names: list[str], timeout: int = 15) -> dict[str, Any]:
    before_list = applescript_list(before_names)
    escaped_prefix = name_prefix.replace("\\", "\\\\").replace('"', '\\"')
    script = f'''
tell application "Keynote"
  set beforeNames to {before_list}
  repeat with i from (count of documents) to 1 by -1
    set d to document i
    set docName to name of d
    if docName starts with "{escaped_prefix}" or beforeNames does not contain docName then
      close d saving no
    end if
  end repeat
end tell
'''
    first = run_osascript(script, timeout)
    time.sleep(2)
    second = run_osascript(script, timeout)
    return {
        "status": second.get("status"),
        "first": first,
        "second": second,
        "remaining_documents": keynote_document_names(timeout),
    }


def export_keynote09(pptx: Path, output: Path, timeout: int) -> dict[str, Any]:
    if output.exists():
        if output.is_dir():
            shutil.rmtree(output)
        else:
            output.unlink()
    script = f'''
tell application "Keynote"
  set inputFile to POSIX file "{applescript_path(pptx.resolve())}"
  set outputFile to POSIX file "{applescript_path(output.resolve())}"
  set docRef to open inputFile
  delay 1
  export docRef to outputFile as Keynote 09
  close docRef saving no
end tell
'''
    result = run_osascript(script, max(timeout, 15))
    result["output_exists"] = output.exists()
    if output.exists():
        result["output_mtime"] = artifact_mtime(output)
    return result


def artifact_mtime(path: Path) -> float:
    if path.is_dir():
        newest = path.stat().st_mtime
        for child in path.rglob("*"):
            try:
                newest = max(newest, child.stat().st_mtime)
            except OSError:
                continue
        return newest
    return path.stat().st_mtime


def create_control_pptx(path: Path) -> dict[str, Any]:
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches, Pt
    except Exception as exc:
        return {"status": STATUS_FAILED, "reason": f"python-pptx unavailable: {exc}"}

    path.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(247, 243, 234)

    title = slide.shapes.add_textbox(Inches(0.85), Inches(0.8), Inches(11.6), Inches(1.0))
    frame = title.text_frame
    frame.text = "Keynote control deck"
    p = frame.paragraphs[0]
    p.font.name = "Aptos Display"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(22, 31, 45)

    body = slide.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(6.2), Inches(2.1))
    body.text_frame.text = "One slide. Editable text. One native rectangle. No SVG, images, charts, or speaker notes."
    body.text_frame.paragraphs[0].font.size = Pt(22)
    body.text_frame.paragraphs[0].font.color.rgb = RGBColor(65, 72, 86)

    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.05), Inches(2.05), Inches(3.7), Inches(2.2))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(200, 90, 46)
    rect.line.color.rgb = RGBColor(200, 90, 46)
    rect.text_frame.text = "native shape"
    rect.text_frame.paragraphs[0].font.size = Pt(24)
    rect.text_frame.paragraphs[0].font.bold = True
    rect.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    presentation.save(path)
    return {"status": STATUS_OK, "path": str(path), "size_bytes": path.stat().st_size}


def build_probe(
    pptx: Path,
    output: Path,
    open_timeout: int,
    save_timeout: int,
    cleanup: bool,
) -> dict[str, Any]:
    report: dict[str, Any] = {
        "schema_version": "1.0.0",
        "tool": "qiaomu-ppt/scripts/keynote_probe.py",
        "generated_at": datetime.now(timezone.utc).isoformat(),
        "pptx": str(pptx),
        "output": str(output),
        "platform": platform.platform(),
        "stages": {},
    }
    stages = report["stages"]

    if platform.system() != "Darwin":
        report["status"] = STATUS_SKIPPED
        report["reason"] = "Keynote probing is macOS-only"
        return report
    if not shutil.which("osascript"):
        report["status"] = STATUS_SKIPPED
        report["reason"] = "osascript not found"
        return report
    keynote_app = Path("/Applications/Keynote.app")
    if not keynote_app.exists():
        report["status"] = STATUS_SKIPPED
        report["reason"] = "Keynote.app not found"
        return report
    if not pptx.exists():
        report["status"] = STATUS_FAILED
        report["reason"] = "PPTX input does not exist"
        return report

    output.parent.mkdir(parents=True, exist_ok=True)
    if output.exists():
        if output.is_dir():
            shutil.rmtree(output)
        else:
            output.unlink()

    before_docs = keynote_document_names()
    report["documents_before"] = before_docs

    stages["version"] = run_osascript('tell application "Keynote" to get version', 15)
    stages["activate"] = run_osascript('tell application "Keynote" to activate', 15)

    open_script = f'''
tell application "Keynote"
  set inputFile to POSIX file "{applescript_path(pptx.resolve())}"
  set docRef to open inputFile
  delay 1
  get name of docRef
end tell
'''
    stages["open"] = run_osascript(open_script, open_timeout)
    doc_name = str(stages["open"].get("stdout") or "").strip().splitlines()[-1:] or [""]
    doc_name = doc_name[0].strip()
    report["opened_document_name"] = doc_name

    if stages["open"].get("status") == STATUS_OK and doc_name:
        save_script = f'''
tell application "Keynote"
  set outputFile to POSIX file "{applescript_path(output.resolve())}"
  set docRef to document {applescript_string(doc_name)}
  save docRef in outputFile as Keynote
end tell
'''
        stages["save"] = run_osascript(save_script, save_timeout)
    else:
        stages["save"] = {
            "status": STATUS_SKIPPED,
            "reason": "open stage did not return a document name",
        }

    if stages["save"].get("status") != STATUS_OK or not output.exists():
        stages["cleanup_after_save_failure"] = cleanup_documents(pptx.stem, before_docs)
        stages["export_keynote09"] = export_keynote09(pptx, output, save_timeout)
    else:
        stages["export_keynote09"] = {"status": STATUS_SKIPPED, "reason": "modern Keynote save produced an output"}

    if cleanup:
        stages["cleanup"] = cleanup_documents(pptx.stem, before_docs)
    else:
        stages["cleanup"] = {"status": STATUS_SKIPPED, "reason": "cleanup disabled"}

    output_exists = output.exists()
    report["output_exists"] = output_exists
    if output_exists:
        report["output_mtime"] = artifact_mtime(output)
        report["pptx_mtime"] = artifact_mtime(pptx)
        report["output_fresh"] = artifact_mtime(output) >= artifact_mtime(pptx)
    else:
        report["output_fresh"] = False

    if output_exists and report["output_fresh"]:
        report["status"] = STATUS_OK
        if stages.get("export_keynote09", {}).get("status") == STATUS_OK:
            report["export_strategy"] = "Keynote 09 fallback"
        else:
            report["export_strategy"] = "save as Keynote"
    elif stages["open"].get("status") == STATUS_TIMEOUT:
        report["status"] = STATUS_TIMEOUT
        report["reason"] = f"Keynote open stage timed out after {open_timeout}s"
    elif stages["save"].get("status") == STATUS_TIMEOUT:
        report["status"] = STATUS_TIMEOUT
        report["reason"] = f"Keynote save stage timed out after {save_timeout}s"
    else:
        report["status"] = STATUS_FAILED
        report["reason"] = "Keynote probe did not produce a fresh .key artifact"
    return report


def build_probe_with_control(
    pptx: Path,
    output: Path,
    report_path: Path,
    open_timeout: int,
    save_timeout: int,
    cleanup: bool,
) -> dict[str, Any]:
    report = build_probe(
        pptx=pptx,
        output=output,
        open_timeout=open_timeout,
        save_timeout=save_timeout,
        cleanup=cleanup,
    )
    control_pptx = report_path.parent / "keynote-control.pptx"
    control_output = report_path.parent / "keynote-control.key"
    control_create = create_control_pptx(control_pptx)
    if control_create.get("status") == STATUS_OK:
        control_probe = build_probe(
            pptx=control_pptx,
            output=control_output,
            open_timeout=open_timeout,
            save_timeout=save_timeout,
            cleanup=cleanup,
        )
    else:
        control_probe = {
            "status": STATUS_FAILED,
            "reason": control_create.get("reason", "control PPTX creation failed"),
        }
    report["control_probe"] = {
        "description": "Minimal one-slide native PowerPoint deck used to distinguish complex-deck issues from baseline Keynote automation issues.",
        "create": control_create,
        "probe": control_probe,
    }
    main_status = report.get("status")
    control_status = control_probe.get("status")
    if main_status == STATUS_OK and report.get("export_strategy") == "Keynote 09 fallback":
        report["diagnosis"] = "input_pptx_keynote09_fallback_ok"
    elif main_status == STATUS_OK:
        report["diagnosis"] = "input_pptx_keynote_save_ok"
    elif control_status == STATUS_OK:
        report["diagnosis"] = "input_pptx_specific_keynote_save_failure"
    elif control_status in {STATUS_TIMEOUT, STATUS_FAILED}:
        report["diagnosis"] = "baseline_keynote_save_automation_failure"
    else:
        report["diagnosis"] = "keynote_probe_inconclusive"
    return report


def main() -> int:
    parser = argparse.ArgumentParser(description="Diagnose Keynote PPTX import/save automation.")
    parser.add_argument("pptx", type=Path, help="PPTX input to import into Keynote.")
    parser.add_argument("--output", type=Path, help="Output .key path. Defaults to sibling .key under the PPTX directory.")
    parser.add_argument("--report", type=Path, help="Diagnostic JSON path. Defaults to <output>.probe.json.")
    parser.add_argument("--open-timeout", type=int, default=120, help="Timeout for Keynote PPTX open/import stage.")
    parser.add_argument("--save-timeout", type=int, default=120, help="Timeout for Keynote save-as-Keynote stage.")
    parser.add_argument("--with-control", action="store_true", help="Also generate and probe a minimal native PPTX control deck.")
    parser.add_argument("--no-cleanup", action="store_true", help="Leave opened Keynote documents in place for manual inspection.")
    parser.add_argument("--strict", action="store_true", help="Exit non-zero when a fresh .key artifact is not produced.")
    args = parser.parse_args()

    pptx = args.pptx.resolve()
    output = args.output.resolve() if args.output else pptx.with_suffix(".key")
    report_path = args.report.resolve() if args.report else output.with_suffix(output.suffix + ".probe.json")
    if args.with_control:
        report = build_probe_with_control(
            pptx=pptx,
            output=output,
            report_path=report_path,
            open_timeout=args.open_timeout,
            save_timeout=args.save_timeout,
            cleanup=not args.no_cleanup,
        )
    else:
        report = build_probe(
            pptx=pptx,
            output=output,
            open_timeout=args.open_timeout,
            save_timeout=args.save_timeout,
            cleanup=not args.no_cleanup,
        )
    write_json(report_path, report)
    print(json.dumps(report, ensure_ascii=False, indent=2))
    if args.strict and report.get("status") != STATUS_OK:
        return 2
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

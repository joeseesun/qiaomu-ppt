#!/usr/bin/env python3
"""Inspect visible text in a generated qiaomu-ppt PPTX."""

from __future__ import annotations

import argparse
import html
import json
import math
import re
import zipfile
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


PLACEHOLDERS = [
    r"\[必填\]",
    r"lorem ipsum",
    r"\bTODO\b",
    r"\bTBD\b",
    r"SLIDES_HERE",
]

FORBIDDEN_VISIBLE = [
    ("fetched via", r"\bfetched\s+via\b"),
    ("generated with", r"\bgenerated\s+with\b"),
    ("qiaomu-markdown-proxy", r"qiaomu-markdown-proxy"),
    ("speaker cue", r"Speaker\s+cue\s*:"),
    ("source footer", r"来源\s*[:：]"),
    ("source_fetch", r"\bsource_fetch\b"),
    ("pptx export", r"\bPPTX\s+export\b"),
    ("presenter framing", r"(讲述角度|叙事角度|设计说明|本页意图|可讲述为)\s*[:：]"),
    ("AI generation wording", r"AI\s*(背景|图像|图片|生成|素材)"),
    ("inspectability wording", r"(作为可检查对象|可检查对象|可检查|可被观察|可观察对象|可被核验|可核验对象)"),
    (
        "real-asset backstage wording",
        r"真实(专辑)?(封面|素材|资料|图片|图像|对象|资产|前景|证据)",
    ),
    (
        "object-authenticity backstage wording",
        r"(封面|素材|资料|图片|图像|对象|资产|前景|证据).{0,8}(真实|可检查|可观察|可被观察|可核验|可被核验)",
    ),
    (
        "foreground/background contract wording",
        r"(前景对象|前景素材|前景层|前景图层|前景证据|背景只(?:提供|承担)|背景(?:提供|承担).{0,12}(气氛|氛围|空间|环境)|背景作为.{0,8}(环境|气氛|氛围))",
    ),
    ("asset placement method", r"(不是.{0,12}(装饰|贴上去)|贴上去|摆放环境)"),
    ("editable-object wording", r"可编辑.{0,8}(对象|文本|图形|前景|PPTX?|PowerPoint)"),
    ("preview/test wording", r"(预览方向|生成测试页|测试页|四页预览|视觉方向验证)"),
    ("toolchain name", r"(ppt-master|qiaomu-ppt|Codex|gpt-image)"),
    ("deck", r"\bdeck\b"),
    ("route", r"\broute\b"),
    ("fallback", r"\bfallback\b"),
    ("artifact", r"\bartifact\b"),
    ("pipeline", r"\bpipeline\b"),
    ("source asset", r"\bsource\s+asset\b"),
    ("foreground layer", r"\bforeground\s+layer\b"),
    ("integration modifier", r"\bintegration\s+modifier"),
    ("background duty", r"\bbackground\s+duty\b"),
    ("image generation queue", r"\bimage\s+generation\s+queue\b"),
    ("sidecar/manifest", r"\b(sidecar|manifest|visual_asset)\b"),
]

EMU_MARGIN_RATIO = 0.06
FULL_SLIDE_COVER_RATIO = 0.88
EMU_PER_INCH = 914400
PT_PER_INCH = 72
TITLE_TOP_REGION_RATIO = 0.36
TEXT_OVERFLOW_FAIL_RATIO = 1.18
TEXT_OVERFLOW_WARN_RATIO = 1.05
TITLE_IMAGE_CLEARANCE_INCH = 0.18
TITLE_CONTENT_FAIL_CLEARANCE_INCH = 0.18
TITLE_CONTENT_WARN_CLEARANCE_INCH = 0.28
MIN_TITLE_TEXT_LEN = 2
ACCENT_CALLOUT_MIN_HEIGHT_INCH = 0.34
ACCENT_INFO_BAR_MIN_HEIGHT_INCH = 0.40
ACCENT_CALLOUT_MAX_ASPECT_RATIO = 16.0


def load_json(path: Path) -> Any:
    with path.open("r", encoding="utf-8") as fh:
        return json.load(fh)


def write_json(path: Path, payload: Any) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(payload, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def iter_slides(plan: Any) -> list[dict[str, Any]]:
    if isinstance(plan, list):
        return [item for item in plan if isinstance(item, dict)]
    if isinstance(plan, dict):
        for key in ("slides", "slide_plan", "pages"):
            value = plan.get(key)
            if isinstance(value, list):
                return [item for item in value if isinstance(item, dict)]
    return []


def expected_title(slide: dict[str, Any]) -> str:
    return str(slide.get("claim_title") or slide.get("title") or "").strip()


def visible_text_by_slide(pptx: Path) -> list[str]:
    presentation = Presentation(str(pptx))
    result: list[str] = []
    for slide in presentation.slides:
        chunks: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text:
                chunks.append(text)
        result.append("\n".join(chunks))
    return result


def visible_text_by_slide_xml(pptx: Path) -> list[str]:
    """Fallback reader for native DrawingML generated outside python-pptx.

    Some SVG-native exporters inject valid DrawingML that is visible in
    PowerPoint/LibreOffice but is not surfaced by python-pptx shape.text.
    Reading the raw slide XML keeps text QA aligned with rendered output.
    """
    def slide_index(name: str) -> int:
        match = re.search(r"slide(\d+)\.xml$", name)
        return int(match.group(1)) if match else 0

    result: list[str] = []
    with zipfile.ZipFile(pptx) as zf:
        names = sorted(
            [name for name in zf.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml")],
            key=slide_index,
        )
        for name in names:
            xml = zf.read(name).decode("utf-8", errors="ignore")
            chunks = [html.unescape(value) for value in re.findall(r"<a:t>(.*?)</a:t>", xml, flags=re.DOTALL)]
            result.append("\n".join(chunks))
    return result


def inspect_editability(pptx: Path) -> dict[str, Any]:
    presentation = Presentation(str(pptx))
    slide_w = int(presentation.slide_width)
    slide_h = int(presentation.slide_height)
    slides: list[dict[str, Any]] = []
    image_backed_count = 0
    full_slide_picture_count = 0
    native_text_chars_total = 0
    native_text_shape_count_total = 0

    for idx, slide in enumerate(presentation.slides, start=1):
        picture_count = 0
        full_slide_pictures = 0
        native_text_chars = 0
        native_text_shapes = 0
        native_non_picture_shapes = 0
        for shape in slide.shapes:
            text = getattr(shape, "text", "") or ""
            if text.strip():
                native_text_shapes += 1
                native_text_chars += len(text.strip())
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                picture_count += 1
                left = int(getattr(shape, "left", 0))
                top = int(getattr(shape, "top", 0))
                width = int(getattr(shape, "width", 0))
                height = int(getattr(shape, "height", 0))
                covers_width = width >= slide_w * FULL_SLIDE_COVER_RATIO and left <= slide_w * EMU_MARGIN_RATIO
                covers_height = height >= slide_h * FULL_SLIDE_COVER_RATIO and top <= slide_h * EMU_MARGIN_RATIO
                if covers_width and covers_height:
                    full_slide_pictures += 1
            else:
                native_non_picture_shapes += 1

        image_backed = full_slide_pictures > 0 and native_text_chars < 12
        if image_backed:
            image_backed_count += 1
        if full_slide_pictures:
            full_slide_picture_count += full_slide_pictures
        native_text_chars_total += native_text_chars
        native_text_shape_count_total += native_text_shapes
        slides.append(
            {
                "slide_no": idx,
                "picture_count": picture_count,
                "full_slide_picture_count": full_slide_pictures,
                "native_text_shape_count": native_text_shapes,
                "native_text_chars": native_text_chars,
                "native_non_picture_shape_count": native_non_picture_shapes,
                "image_backed": image_backed,
            }
        )

    slide_count = len(slides)
    return {
        "slide_count": slide_count,
        "native_text_shape_count": native_text_shape_count_total,
        "native_text_chars": native_text_chars_total,
        "full_slide_picture_count": full_slide_picture_count,
        "image_backed_slide_count": image_backed_count,
        "image_backed_ratio": image_backed_count / slide_count if slide_count else 0,
        "slides": slides,
    }


def emu_to_pt(value: int | float) -> float:
    return float(value) / EMU_PER_INCH * PT_PER_INCH


def shape_bounds(shape: Any) -> dict[str, int]:
    left = int(getattr(shape, "left", 0) or 0)
    top = int(getattr(shape, "top", 0) or 0)
    width = max(0, int(getattr(shape, "width", 0) or 0))
    height = max(0, int(getattr(shape, "height", 0) or 0))
    return {
        "left": left,
        "top": top,
        "width": width,
        "height": height,
        "right": left + width,
        "bottom": top + height,
    }


def shape_fill_rgb(shape: Any) -> tuple[int, int, int] | None:
    try:
        fill = getattr(shape, "fill", None)
        if fill is None:
            return None
        fore_color = fill.fore_color
        rgb = fore_color.rgb
        if rgb is None:
            return None
        return int(rgb[0]), int(rgb[1]), int(rgb[2])
    except Exception:
        return None


def is_accent_fill(rgb: tuple[int, int, int] | None) -> bool:
    if rgb is None:
        return False
    red, green, blue = rgb
    saturated_red = red >= 150 and green <= 115 and blue <= 115
    saturated_blue = blue >= 140 and red <= 150 and green >= 90
    saturated_green = green >= 140 and red <= 150 and blue <= 150
    return saturated_red or saturated_blue or saturated_green


def is_rounded_callout_shape(shape: Any) -> bool:
    auto_shape_type = getattr(shape, "auto_shape_type", None)
    if auto_shape_type is None:
        return False
    name = str(auto_shape_type).upper()
    return "ROUND" in name


def inspect_accent_callout(
    *,
    slide_no: int,
    shape_index: int,
    shape: Any,
    text: str,
    bounds: dict[str, int],
) -> tuple[list[str], list[str], dict[str, Any] | None]:
    rgb = shape_fill_rgb(shape)
    if not is_accent_fill(rgb):
        return [], [], None
    normalized = re.sub(r"\s+", "", text)
    if not normalized:
        return [], [], None
    width_in = bounds["width"] / EMU_PER_INCH
    height_in = bounds["height"] / EMU_PER_INCH
    aspect_ratio = width_in / height_in if height_in else 999.0
    report = {
        "shape_index": shape_index,
        "text": re.sub(r"\s+", " ", text)[:80],
        "width_inch": round(width_in, 3),
        "height_inch": round(height_in, 3),
        "aspect_ratio": round(aspect_ratio, 2),
        "rounded": is_rounded_callout_shape(shape),
        "fill_rgb": list(rgb) if rgb else None,
    }
    failures: list[str] = []
    warnings: list[str] = []
    if not report["rounded"]:
        failures.append(
            f"slide {slide_no} accent callout should be a rounded rectangle: "
            f"shape {shape_index} text={report['text']}"
        )
    if height_in < ACCENT_CALLOUT_MIN_HEIGHT_INCH:
        failures.append(
            f"slide {slide_no} accent callout is too short: shape {shape_index} "
            f"height={height_in:.2f}in text={report['text']}"
        )
    if len(normalized) >= 8 and height_in < ACCENT_INFO_BAR_MIN_HEIGHT_INCH:
        failures.append(
            f"slide {slide_no} accent info bar needs more vertical breathing room: "
            f"shape {shape_index} height={height_in:.2f}in text={report['text']}"
        )
    if aspect_ratio > ACCENT_CALLOUT_MAX_ASPECT_RATIO:
        failures.append(
            f"slide {slide_no} accent callout is too skinny: shape {shape_index} "
            f"aspect={aspect_ratio:.1f} text={report['text']}"
        )
    elif aspect_ratio > 12 and height_in < 0.48:
        warnings.append(
            f"slide {slide_no} accent callout is wide; verify it reads as an info bar, "
            f"not a flattened pill: shape {shape_index} aspect={aspect_ratio:.1f} text={report['text']}"
        )
    return failures, warnings, report


def rect_area(rect: dict[str, int]) -> int:
    return max(0, rect["width"]) * max(0, rect["height"])


def rect_intersection(a: dict[str, int], b: dict[str, int]) -> dict[str, int] | None:
    left = max(a["left"], b["left"])
    top = max(a["top"], b["top"])
    right = min(a["right"], b["right"])
    bottom = min(a["bottom"], b["bottom"])
    if right <= left or bottom <= top:
        return None
    return {"left": left, "top": top, "right": right, "bottom": bottom, "width": right - left, "height": bottom - top}


def horizontal_overlap(a: dict[str, int], b: dict[str, int]) -> int:
    return max(0, min(a["right"], b["right"]) - max(a["left"], b["left"]))


def vertical_gap(a: dict[str, int], b: dict[str, int]) -> int | None:
    if b["top"] >= a["bottom"]:
        return b["top"] - a["bottom"]
    if a["top"] >= b["bottom"]:
        return a["top"] - b["bottom"]
    return None


def expand_rect(rect: dict[str, int], amount: int) -> dict[str, int]:
    return {
        "left": rect["left"] - amount,
        "top": rect["top"] - amount,
        "right": rect["right"] + amount,
        "bottom": rect["bottom"] + amount,
        "width": rect["width"] + amount * 2,
        "height": rect["height"] + amount * 2,
    }


def text_frame_margins_pt(shape: Any) -> tuple[float, float, float, float]:
    text_frame = getattr(shape, "text_frame", None)
    if text_frame is None:
        return (7.2, 7.2, 3.6, 3.6)
    left = emu_to_pt(int(getattr(text_frame, "margin_left", 0) or 0))
    right = emu_to_pt(int(getattr(text_frame, "margin_right", 0) or 0))
    top = emu_to_pt(int(getattr(text_frame, "margin_top", 0) or 0))
    bottom = emu_to_pt(int(getattr(text_frame, "margin_bottom", 0) or 0))
    return (left, right, top, bottom)


def text_shape_max_font_pt(shape: Any) -> float:
    text_frame = getattr(shape, "text_frame", None)
    max_size = 0.0
    if text_frame is not None:
        for paragraph in text_frame.paragraphs:
            p_size = getattr(getattr(paragraph, "font", None), "size", None)
            if p_size:
                max_size = max(max_size, float(p_size.pt))
            for run in paragraph.runs:
                r_size = getattr(getattr(run, "font", None), "size", None)
                if r_size:
                    max_size = max(max_size, float(r_size.pt))
    if max_size:
        return max_size
    height_pt = emu_to_pt(int(getattr(shape, "height", 0) or 0))
    return max(11.0, min(28.0, height_pt * 0.38))


def weighted_line_units(text: str) -> float:
    units = 0.0
    for char in text:
        code = ord(char)
        if char.isspace():
            units += 0.28
        elif 0x2E80 <= code <= 0x9FFF or 0xFF00 <= code <= 0xFFEF:
            units += 1.0
        elif char.isupper():
            units += 0.68
        elif char.isalnum():
            units += 0.56
        else:
            units += 0.46
    return units


def estimate_text_fit(shape: Any) -> dict[str, Any] | None:
    text = (getattr(shape, "text", "") or "").strip()
    if not text:
        return None
    bounds = shape_bounds(shape)
    if bounds["width"] <= 0 or bounds["height"] <= 0:
        return None
    font_pt = text_shape_max_font_pt(shape)
    margin_left, margin_right, margin_top, margin_bottom = text_frame_margins_pt(shape)
    available_width_pt = max(8.0, emu_to_pt(bounds["width"]) - margin_left - margin_right)
    available_height_pt = max(4.0, emu_to_pt(bounds["height"]) - margin_top - margin_bottom)
    units_per_line = max(1.0, available_width_pt / max(1.0, font_pt * 0.86))
    text_frame = getattr(shape, "text_frame", None)
    paragraph_texts = [text]
    paragraph_count = 1
    if text_frame is not None:
        paragraph_texts = [paragraph.text for paragraph in text_frame.paragraphs]
        paragraph_count = len(paragraph_texts)
    required_lines = 0.0
    for paragraph_text in paragraph_texts:
        explicit_lines = paragraph_text.splitlines() or [paragraph_text]
        for line in explicit_lines:
            units = weighted_line_units(line.strip())
            required_lines += max(1, math.ceil(units / units_per_line)) if units else 0.35
    required_height_pt = required_lines * font_pt * 1.18 + max(0, paragraph_count - 1) * font_pt * 0.2
    ratio = required_height_pt / available_height_pt
    return {
        "text_preview": re.sub(r"\s+", " ", text)[:80],
        "font_pt": round(font_pt, 2),
        "required_lines": round(required_lines, 2),
        "available_width_pt": round(available_width_pt, 2),
        "available_height_pt": round(available_height_pt, 2),
        "required_height_pt": round(required_height_pt, 2),
        "height_ratio": round(ratio, 3),
        "bounds": bounds,
    }


def is_background_picture(rect: dict[str, int], slide_w: int, slide_h: int) -> bool:
    slide_area = slide_w * slide_h
    if slide_area <= 0:
        return False
    area_ratio = rect_area(rect) / slide_area
    near_full_width = rect["width"] >= slide_w * 0.82 and rect["left"] <= slide_w * 0.09
    near_full_height = rect["height"] >= slide_h * 0.82 and rect["top"] <= slide_h * 0.09
    return area_ratio >= 0.72 and (near_full_width or near_full_height)


def text_item_is_footer(item: dict[str, Any], slide_h: int) -> bool:
    text = item["text"].strip()
    bounds = item["bounds"]
    if bounds["top"] >= slide_h * 0.82:
        return True
    if re.fullmatch(r"\d{1,2}\s*/\s*.+", text):
        return True
    if re.search(r"\b(source|copyright|github\.com|https?://)\b", text, flags=re.IGNORECASE):
        return True
    return False


def text_item_is_minor_label(item: dict[str, Any]) -> bool:
    text = item["text"].strip()
    if len(text) <= 2 and item["font_pt"] <= 18:
        return True
    return False


def title_candidates_for_slide(text_items: list[dict[str, Any]], slide_h: int) -> list[dict[str, Any]]:
    if not text_items:
        return []
    max_font = max(item["font_pt"] for item in text_items)
    return [
        item
        for item in text_items
        if item["bounds"]["top"] <= slide_h * TITLE_TOP_REGION_RATIO
        and item["font_pt"] >= max(18.0, max_font * 0.68)
        and len(item["text"]) >= MIN_TITLE_TEXT_LEN
        and not text_item_is_footer(item, slide_h)
    ]


def inspect_title_spacing(
    *,
    slide_no: int,
    title: dict[str, Any],
    text_items: list[dict[str, Any]],
    fail_gap: int,
    warn_gap: int,
    slide_h: int,
) -> tuple[list[str], list[str], list[dict[str, Any]]]:
    failures: list[str] = []
    warnings: list[str] = []
    contacts: list[dict[str, Any]] = []
    title_rect = title["bounds"]
    for item in text_items:
        if item["shape_index"] == title["shape_index"]:
            continue
        if text_item_is_footer(item, slide_h) or text_item_is_minor_label(item):
            continue
        item_rect = item["bounds"]
        if item_rect["top"] < title_rect["bottom"]:
            continue
        x_overlap = horizontal_overlap(title_rect, item_rect)
        min_width = max(1, min(title_rect["width"], item_rect["width"]))
        if x_overlap / min_width < 0.18:
            continue
        gap = vertical_gap(title_rect, item_rect)
        if gap is None:
            continue
        gap_in = gap / EMU_PER_INCH
        contact = {
            "title_shape_index": title["shape_index"],
            "content_shape_index": item["shape_index"],
            "title_text": title["text"],
            "content_text": item["text"],
            "gap_inch": round(gap_in, 3),
        }
        if gap < fail_gap:
            contacts.append(contact)
            failures.append(
                f"slide {slide_no} title is too close to following content: "
                f"title shape {title['shape_index']} content shape {item['shape_index']} "
                f"gap={gap_in:.2f}in title={title['text']} content={item['text']}"
            )
        elif gap < warn_gap:
            contacts.append(contact)
            warnings.append(
                f"slide {slide_no} title breathing room is tight: "
                f"title shape {title['shape_index']} content shape {item['shape_index']} "
                f"gap={gap_in:.2f}in title={title['text']} content={item['text']}"
            )
    return failures, warnings, contacts


def inspect_layout(pptx: Path) -> dict[str, Any]:
    presentation = Presentation(str(pptx))
    slide_w = int(presentation.slide_width)
    slide_h = int(presentation.slide_height)
    slide_area = slide_w * slide_h
    image_clearance = int(TITLE_IMAGE_CLEARANCE_INCH * EMU_PER_INCH)
    content_fail_gap = int(TITLE_CONTENT_FAIL_CLEARANCE_INCH * EMU_PER_INCH)
    content_warn_gap = int(TITLE_CONTENT_WARN_CLEARANCE_INCH * EMU_PER_INCH)
    failures: list[str] = []
    warnings: list[str] = []
    slides: list[dict[str, Any]] = []

    for idx, slide in enumerate(presentation.slides, start=1):
        text_items: list[dict[str, Any]] = []
        picture_items: list[dict[str, Any]] = []
        slide_report: dict[str, Any] = {
            "slide_no": idx,
            "text_fit": [],
            "accent_callouts": [],
            "title_image_contacts": [],
            "title_content_contacts": [],
        }
        for shape_index, shape in enumerate(slide.shapes, start=1):
            bounds = shape_bounds(shape)
            text = (getattr(shape, "text", "") or "").strip()
            if text:
                accent_failures, accent_warnings, accent_report = inspect_accent_callout(
                    slide_no=idx,
                    shape_index=shape_index,
                    shape=shape,
                    text=text,
                    bounds=bounds,
                )
                failures.extend(accent_failures)
                warnings.extend(accent_warnings)
                if accent_report:
                    slide_report["accent_callouts"].append(accent_report)
                fit = estimate_text_fit(shape)
                if fit:
                    fit["shape_index"] = shape_index
                    slide_report["text_fit"].append(fit)
                    if fit["height_ratio"] > TEXT_OVERFLOW_FAIL_RATIO:
                        failures.append(
                            f"slide {idx} text likely overflows shape {shape_index}: "
                            f"ratio={fit['height_ratio']} text={fit['text_preview']}"
                        )
                    elif fit["height_ratio"] > TEXT_OVERFLOW_WARN_RATIO and len(fit["text_preview"].strip()) > 2:
                        warnings.append(
                            f"slide {idx} text is tight in shape {shape_index}: "
                            f"ratio={fit['height_ratio']} text={fit['text_preview']}"
                        )
                text_items.append(
                    {
                        "shape_index": shape_index,
                        "bounds": bounds,
                        "font_pt": text_shape_max_font_pt(shape),
                        "text": re.sub(r"\s+", " ", text)[:80],
                    }
                )
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                picture_items.append({"shape_index": shape_index, "bounds": bounds})

        title_candidates = title_candidates_for_slide(text_items, slide_h)
        for title in title_candidates:
            title_rect = title["bounds"]
            spacing_failures, spacing_warnings, spacing_contacts = inspect_title_spacing(
                slide_no=idx,
                title=title,
                text_items=text_items,
                fail_gap=content_fail_gap,
                warn_gap=content_warn_gap,
                slide_h=slide_h,
            )
            failures.extend(spacing_failures)
            warnings.extend(spacing_warnings)
            slide_report["title_content_contacts"].extend(spacing_contacts)
            protected_title_rect = expand_rect(title_rect, image_clearance)
            for picture in picture_items:
                picture_rect = picture["bounds"]
                if slide_area and rect_area(picture_rect) / slide_area < 0.025:
                    continue
                if is_background_picture(picture_rect, slide_w, slide_h):
                    continue
                overlap = rect_intersection(title_rect, picture_rect)
                clearance_hit = rect_intersection(protected_title_rect, picture_rect)
                if overlap:
                    contact = {
                        "title_shape_index": title["shape_index"],
                        "picture_shape_index": picture["shape_index"],
                        "title_text": title["text"],
                        "overlap": overlap,
                    }
                    slide_report["title_image_contacts"].append(contact)
                    failures.append(
                        f"slide {idx} title overlaps image: title shape {title['shape_index']} "
                        f"image shape {picture['shape_index']} text={title['text']}"
                    )
                elif clearance_hit:
                    contact = {
                        "title_shape_index": title["shape_index"],
                        "picture_shape_index": picture["shape_index"],
                        "title_text": title["text"],
                        "clearance_inch": TITLE_IMAGE_CLEARANCE_INCH,
                    }
                    slide_report["title_image_contacts"].append(contact)
                    failures.append(
                        f"slide {idx} title is too close to image: title shape {title['shape_index']} "
                        f"image shape {picture['shape_index']} text={title['text']}"
                    )
        slides.append(slide_report)

    return {
        "text_overflow_fail_ratio": TEXT_OVERFLOW_FAIL_RATIO,
        "text_overflow_warn_ratio": TEXT_OVERFLOW_WARN_RATIO,
        "title_image_clearance_inch": TITLE_IMAGE_CLEARANCE_INCH,
        "title_content_fail_clearance_inch": TITLE_CONTENT_FAIL_CLEARANCE_INCH,
        "title_content_warn_clearance_inch": TITLE_CONTENT_WARN_CLEARANCE_INCH,
        "failures": failures,
        "warnings": warnings,
        "slides": slides,
    }


def best_visible_text_by_slide(pptx: Path) -> tuple[list[str], str]:
    primary = visible_text_by_slide(pptx)
    if sum(len(text.strip()) for text in primary) >= max(1, len(primary)) * 12:
        return primary, "python-pptx"
    fallback = visible_text_by_slide_xml(pptx)
    return fallback, "pptx-slide-xml"


def normalize_text(value: str) -> str:
    value = re.sub(r"\s+", "", value)
    return re.sub(r"[，,。.:：；;、＋+\-—–·/／()（）\[\]【】\"“”'‘’]", "", value)


def scan_patterns(text: str, patterns: list[str]) -> list[str]:
    return [pattern for pattern in patterns if re.search(pattern, text, flags=re.IGNORECASE)]


def scan_labeled_patterns(text: str, patterns: list[tuple[str, str]]) -> list[str]:
    return [label for label, pattern in patterns if re.search(pattern, text, flags=re.IGNORECASE)]


def check(pptx: Path, slide_plan: Path | None, *, allow_image_backed: bool = False) -> dict[str, Any]:
    texts, extraction_method = best_visible_text_by_slide(pptx)
    editability = inspect_editability(pptx)
    layout = inspect_layout(pptx)
    failures: list[str] = []
    warnings: list[str] = []
    expected: list[dict[str, Any]] = []
    if slide_plan and slide_plan.exists():
        expected = iter_slides(load_json(slide_plan))
        if expected and len(expected) != len(texts):
            failures.append(f"slide count mismatch: pptx={len(texts)} slide_plan={len(expected)}")
        for idx, slide in enumerate(expected[: len(texts)], start=1):
            title = expected_title(slide)
            if title and normalize_text(title) not in normalize_text(texts[idx - 1]):
                failures.append(f"slide {idx} missing expected title: {title}")
    for idx, text in enumerate(texts, start=1):
        placeholder_hits = scan_patterns(text, PLACEHOLDERS)
        if placeholder_hits:
            failures.append(f"slide {idx} has placeholder residue: {', '.join(placeholder_hits)}")
        forbidden_hits = scan_labeled_patterns(text, FORBIDDEN_VISIBLE)
        if forbidden_hits:
            failures.append(
                f"slide {idx} has forbidden audience-facing production/backstage text: "
                + ", ".join(forbidden_hits)
            )
        if len(text.strip()) < 12:
            warnings.append(f"slide {idx} has very little visible text")
    image_backed_slides = [
        str(item["slide_no"]) for item in editability["slides"] if item.get("image_backed")
    ]
    if image_backed_slides and not allow_image_backed:
        failures.append(
            "slides appear to be image-backed instead of editable native PPTX content: "
            + ", ".join(image_backed_slides)
        )
    if editability["image_backed_ratio"] >= 0.8 and not allow_image_backed:
        failures.append(
            "PPTX is mostly whole-slide raster images; this is not acceptable for normal editable PPTX delivery"
        )
    elif image_backed_slides and allow_image_backed:
        warnings.append(
            "image-backed slides allowed by flag; label this artifact as non-editable/parity/social-image output"
        )
    failures.extend(layout["failures"])
    warnings.extend(layout["warnings"])
    return {
        "ok": not failures,
        "pptx": str(pptx),
        "extraction_method": extraction_method,
        "slide_count": len(texts),
        "editability": editability,
        "layout": layout,
        "allow_image_backed": allow_image_backed,
        "failures": failures,
        "warnings": warnings,
    }


def main() -> int:
    parser = argparse.ArgumentParser(description="Check visible text in qiaomu-ppt PPTX exports.")
    parser.add_argument("pptx", help="PPTX file to inspect")
    parser.add_argument("--slide-plan", help="Optional slide_plan.json for title/slide-count parity")
    parser.add_argument("--output", "-o", help="Optional JSON report path")
    parser.add_argument(
        "--allow-image-backed",
        action="store_true",
        help="Allow whole-slide raster PPTX exports for explicitly labelled parity/social-image artifacts.",
    )
    args = parser.parse_args()
    result = check(
        Path(args.pptx).resolve(),
        Path(args.slide_plan).resolve() if args.slide_plan else None,
        allow_image_backed=args.allow_image_backed,
    )
    if args.output:
        write_json(Path(args.output), result)
    print(json.dumps(result, ensure_ascii=False, indent=2))
    return 0 if result["ok"] else 2


if __name__ == "__main__":
    raise SystemExit(main())
